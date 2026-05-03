"""Generate a Product Hunt launch presentation for mem-dog."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x60, 0xA5, 0xFA)  # blue-400
ACCENT2 = RGBColor(0x34, 0xD3, 0x99)  # emerald-400
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


# ============================================================
# SLIDE 1: Title / Hero
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
         "mem-dog", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1),
         "The Private AI Memory System for Individuals & Teams",
         font_size=28, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.2), Inches(9), Inches(1.5),
         "Ingest data from 300+ apps. Enrich it with 42 AI agents.\n"
         "Query it with 5 search modes and a temporal knowledge graph.",
         font_size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(3), Inches(6.2), Inches(7), Inches(0.6),
         "Open Source  |  Self-Hosted  |  Apache 2.0",
         font_size=16, color=ACCENT2, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "The Problem", font_size=36, bold=True, color=WHITE)

problems = [
    "Your data is scattered across dozens of apps (Slack, Gmail, Drive, Notion...)",
    "AI assistants forget everything between sessions",
    "Existing memory tools only solve one piece: storage OR search OR integrations",
    "Cloud-only solutions mean your private data lives on someone else's servers",
    "No single system connects ingestion + AI enrichment + knowledge graph + search",
]
add_bullet_list(slide, Inches(1), Inches(2), Inches(11), Inches(5),
                [f"\u2022  {p}" for p in problems], font_size=20, color=GRAY)


# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "mem-dog: One Platform, Complete AI Memory", font_size=32, bold=True, color=WHITE)

solutions = [
    ("Ingest", "Connect 300+ apps via OAuth. Data flows in automatically."),
    ("Enrich", "42 AI agents classify, analyze, and embed every piece of data."),
    ("Store", "Versioned storage with per-item access control and 10 memory types."),
    ("Connect", "Temporal knowledge graph links entities across time."),
    ("Search", "5 search modes + 4 rerankers. RAG chat with citations."),
]

for i, (title, desc) in enumerate(solutions):
    y = Inches(2.0 + i * 1.0)
    add_text(slide, Inches(1.2), y, Inches(2), Inches(0.6),
             title, font_size=22, bold=True, color=ACCENT)
    add_text(slide, Inches(3.5), y, Inches(8.5), Inches(0.6),
             desc, font_size=18, color=GRAY)


# ============================================================
# SLIDE 4: How It Works (Architecture)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         "How It Works", font_size=36, bold=True, color=WHITE)

# Simplified flow
flow_items = [
    ("1. Connect", "OAuth into Slack, Gmail, Drive, Notion, 300+ more"),
    ("2. Ingest", "Webhook Gateway normalizes all data into a universal format"),
    ("3. Process", "42 AI agents classify content, extract entities, generate embeddings"),
    ("4. Store", "Dual-write to Postgres (pgvector) + Neo4j (temporal knowledge graph)"),
    ("5. Query", "Vector, keyword, hybrid, graph, or full search \u2014 with RAG chat"),
]

for i, (step, detail) in enumerate(flow_items):
    y = Inches(1.7 + i * 1.1)
    add_text(slide, Inches(1), y, Inches(3), Inches(0.5),
             step, font_size=20, bold=True, color=ACCENT2)
    add_text(slide, Inches(4), y, Inches(8.5), Inches(0.7),
             detail, font_size=17, color=GRAY)


# ============================================================
# SLIDE 5: Key Differentiators
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         "What Makes mem-dog Different", font_size=36, bold=True, color=WHITE)

diffs = [
    "\u2022  300+ integrations \u2014 not just an API, a full ingestion platform",
    "\u2022  42 typed AI agents \u2014 data-type-specific enrichment, not one-size-fits-all",
    "\u2022  Temporal knowledge graph \u2014 ask \"who was X in 2024?\" and get the right answer",
    "\u2022  5 search modes + reranking \u2014 vector, keyword, hybrid, graph, full",
    "\u2022  Local-first models \u2014 runs offline with Ollama (4b to 27b parameters)",
    "\u2022  Privacy by design \u2014 self-hosted, your data never leaves your infra",
    "\u2022  SDKs in 5 languages \u2014 Python, TypeScript, Go, Rust, Ruby",
    "\u2022  MCP server \u2014 works with Claude Desktop, Cursor, and any MCP client",
]
add_bullet_list(slide, Inches(1), Inches(1.7), Inches(11), Inches(5.5),
                diffs, font_size=19, color=GRAY)


# ============================================================
# SLIDE 6: Competitive Landscape
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8),
         "mem-dog vs. Alternatives", font_size=36, bold=True, color=WHITE)

# Simple comparison text
comparisons = [
    "mem0 \u2014 Memory SDK, vector-only search, no integrations, no knowledge graph",
    "Zep \u2014 Knowledge graph (Graphiti) but no ingestion platform, no UI, limited search",
    "Dify.ai \u2014 LLM workflow builder, no temporal reasoning, no multi-signal search",
    "LangMem \u2014 LangChain plugin, vector-only, no enrichment pipeline",
    "",
    "mem-dog = mem0 + Zep + Nango + Dify in one self-hosted platform",
]

add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5),
                [f"\u2022  {c}" if c else "" for c in comparisons], font_size=19, color=GRAY)

# Highlight the summary line
add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.7),
         "All open source. All self-hosted. ~$50-100/mo with local models.",
         font_size=20, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: Use Cases
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         "Who Is This For?", font_size=36, bold=True, color=WHITE)

cases = [
    ("Individuals", "Personal AI memory across all your apps \u2014 Slack, email, docs, notes"),
    ("Developers", "Add long-term memory to any AI agent with 5 SDKs or MCP"),
    ("Teams", "Shared organizational knowledge with access controls and audit trails"),
    ("Privacy-Conscious", "Self-hosted, local models, your data stays on your hardware"),
    ("AI Builders", "Drop-in memory layer for LangChain, CrewAI, OpenAI agents"),
]

for i, (who, what) in enumerate(cases):
    y = Inches(1.7 + i * 1.1)
    add_text(slide, Inches(1.2), y, Inches(3), Inches(0.6),
             who, font_size=21, bold=True, color=ORANGE)
    add_text(slide, Inches(4.5), y, Inches(8), Inches(0.7),
             what, font_size=17, color=GRAY)


# ============================================================
# SLIDE 8: Tech Stack At a Glance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         "Tech Stack", font_size=36, bold=True, color=WHITE)

stack = [
    "\u2022  API: Python 3.12, FastAPI, 70+ endpoints",
    "\u2022  UI: Next.js 14, TypeScript, 6-mode upload, Knowledge Chat",
    "\u2022  Pipeline: NATS messaging, 42 ADK agents, 5 model tiers (4b\u201327b)",
    "\u2022  Storage: Supabase Postgres + pgvector, Neo4j + Graphiti, GCS",
    "\u2022  Integrations: Nango (300+ OAuth providers, auto token refresh)",
    "\u2022  Search: pgvector cosine, BM25, RRF/MMR/cross-encoder reranking",
    "\u2022  Deploy: Docker Compose (local) or GKE/Cloud Run (production)",
    "\u2022  SDKs: Python, TypeScript, Go, Rust, Ruby + MCP server",
]
add_bullet_list(slide, Inches(1), Inches(1.7), Inches(11), Inches(5.5),
                stack, font_size=19, color=GRAY)


# ============================================================
# SLIDE 9: Quick Start
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Get Started in 60 Seconds", font_size=36, bold=True, color=WHITE)

add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(3.5),
         "git clone https://github.com/mem-dog/mem-dog\n"
         "cd mem-dog\n"
         "docker compose up\n\n"
         "# UI:      http://localhost:3000\n"
         "# API:     http://localhost:8080/docs\n"
         "# Neo4j:   http://localhost:7474",
         font_size=22, color=ACCENT, align=PP_ALIGN.LEFT)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.8),
         "10 services. Zero config. Runs on any machine with Docker.",
         font_size=18, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(1.5), Inches(2), Inches(10), Inches(1.2),
         "Give Your AI a Memory It Deserves",
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1),
         "Open source  \u2022  Self-hosted  \u2022  Apache 2.0",
         font_size=22, color=ACCENT2, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5), Inches(10), Inches(1.5),
         "\u2B50  Star us on GitHub  \u2022  Try it locally  \u2022  Join the community",
         font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.6),
         "github.com/mem-dog/mem-dog",
         font_size=18, color=GRAY, align=PP_ALIGN.CENTER)


# Save
output_path = "/Users/paragagarwal/repos/mem-dog/docs/presentations/mem-dog-product-hunt.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
