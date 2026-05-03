#!/usr/bin/env python3
"""Generate ingestion pipeline + AI enrichment presentation.

Design: one idea per slide, large text, plenty of whitespace.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_HEADER_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_ROW_BG = RGBColor(0x15, 0x1F, 0x32)
TABLE_ALT_BG = RGBColor(0x1A, 0x25, 0x38)
GREEN = RGBColor(0x34, 0xD3, 0x99)
ORANGE = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def slide_title(slide, title):
    add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
             title, size=38, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(1.35), Inches(2.5), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def card(slide, left, top, width, height, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()


def make_table(slide, left, top, width, rows_data, font_size=15):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                   width, Inches(0.45 * n_rows))
    table = shape.table
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Calibri"
                p.font.bold = r == 0
                p.font.color.rgb = ACCENT if r == 0 else WHITE
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = (TABLE_HEADER_BG if r == 0
                                   else TABLE_ALT_BG if r % 2 == 0
                                   else TABLE_ROW_BG)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         "From Ingestion to Intelligence", size=52, bold=True,
         align=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(5), Inches(3.0), Inches(3.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
         "How Mem-Dog Ingests, Enriches, and Makes Data Queryable",
         size=26, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(5.2), Inches(9.333), Inches(0.8),
         "3 ingestion paths  \u00b7  6-layer classification  \u00b7  42 AI agents\n"
         "5 search modes  \u00b7  4 rerankers  \u00b7  RAG chat with citations",
         size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Pipeline Overview
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "The Pipeline at a Glance")

# Single horizontal flow
phases = [
    ("Ingest", "3 entry paths\ninto the system", ACCENT),
    ("Classify", "6-layer detection\nfinds the right agent", ORANGE),
    ("Enrich", "42 AI agents\nanalyze content", GREEN),
    ("Store", "Viewpoints, embeddings,\nentities, graph", ACCENT2),
    ("Query", "5 search modes\nRAG chat + citations", GREEN),
]

box_w = Inches(2.1)
gap = Inches(0.3)
total = len(phases) * box_w + (len(phases) - 1) * gap
start_x = (Inches(13.333) - total) / 2

for i, (title, sub, color) in enumerate(phases):
    left = start_x + i * (box_w + gap)
    card(slide, left, Inches(2.2), box_w, Inches(2.0), border_color=color)
    add_text(slide, left + Inches(0.2), Inches(2.4), box_w - Inches(0.4), Inches(0.6),
             title, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), Inches(3.1), box_w - Inches(0.4), Inches(0.8),
             sub, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < len(phases) - 1:
        add_text(slide, left + box_w, Inches(2.8), gap, Inches(0.5),
                 "\u2192", size=26, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.8),
         "Data enters from any source, gets classified without AI when possible,\n"
         "enriched by a specialized agent, and stored in multiple queryable indexes.",
         size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Three Ingestion Paths
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "Three Ways Data Enters")

paths = [
    ("Channel Webhooks", ACCENT,
     "33 channel adapters in the Webhook Gateway normalize "
     "messages from Slack, WhatsApp, Email, GitHub, Jira, etc. "
     "into a UniversalEnvelope and forward to the API.",
     "POST /webhooks/{webhook_id}"),
    ("Direct API Upload", GREEN,
     "Upload files, text, or URLs directly through the API "
     "or UI. Supports 6 upload modes: text, file, URL, "
     "camera, voice, and video.",
     "POST /api/v1/data"),
    ("Universal Envelope", ORANGE,
     "Structured ingest with origin, payload, and context "
     "descriptors. Supports direct=true to bypass the "
     "pipeline for pre-processed data.",
     "POST /api/v1/ingest"),
]

for i, (title, color, desc, endpoint) in enumerate(paths):
    top = Inches(1.8) + i * Inches(1.8)
    card(slide, Inches(1), top, Inches(11.333), Inches(1.5), border_color=color)
    add_text(slide, Inches(1.5), top + Inches(0.15), Inches(3), Inches(0.5),
             title, size=24, color=color, bold=True)
    add_text(slide, Inches(1.5), top + Inches(0.65), Inches(3), Inches(0.6),
             endpoint, size=15, color=GREEN, font="Courier New")
    add_text(slide, Inches(5), top + Inches(0.25), Inches(7), Inches(1.0),
             desc, size=17, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — UniversalEnvelope
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "The Universal Envelope")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "Every piece of data is wrapped in a UniversalEnvelope \u2014 a single "
         "container that standardizes any data type for the pipeline.",
         size=19, color=LIGHT_GRAY)

# Three sections of the envelope
sections = [
    ("Origin", ACCENT,
     "Where did this come from?",
     "source_type  \u2014  email, slack, api, file...\n"
     "channel_type  \u2014  whatsapp, telegram...\n"
     "user_id  \u2014  who sent it\n"
     "device_id  \u2014  what device\n"
     "timestamp  \u2014  when it happened"),
    ("Payload", GREEN,
     "What is the actual content?",
     "content_type  \u2014  text, file, url, binary\n"
     "data  \u2014  the raw content\n"
     "mime_type  \u2014  application/pdf, etc.\n"
     "metadata  \u2014  tags, labels, context\n"
     "attachments  \u2014  files, images"),
    ("Context", ORANGE,
     "What\u2019s the surrounding context?",
     "memory_type  \u2014  conversation, session...\n"
     "group_id  \u2014  link related items\n"
     "thread_id  \u2014  conversation thread\n"
     "project_id  \u2014  org/project scope\n"
     "tags  \u2014  user-defined labels"),
]

for i, (title, color, question, fields) in enumerate(sections):
    left = Inches(1) + i * Inches(3.9)
    card(slide, left, Inches(2.8), Inches(3.5), Inches(3.8), border_color=color)
    add_text(slide, left + Inches(0.4), Inches(3.0), Inches(2.7), Inches(0.5),
             title, size=24, color=color, bold=True)
    add_text(slide, left + Inches(0.4), Inches(3.5), Inches(2.7), Inches(0.4),
             question, size=15, color=WHITE)
    add_text(slide, left + Inches(0.4), Inches(4.1), Inches(2.7), Inches(2.2),
             fields, size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — NATS + Pull Worker
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "NATS Pipeline Dispatch")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "After the API stores the raw data, it publishes to NATS JetStream. "
         "A pull worker picks up the message and routes it to the right agent.",
         size=19, color=LIGHT_GRAY)

# Left: how it works
card(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(3.8))
add_text(slide, Inches(1.4), Inches(3.0), Inches(4.7), Inches(0.5),
         "Dispatch Flow", size=22, color=ACCENT, bold=True)

steps = [
    ("API stores data in Postgres + blob storage", WHITE),
    ("API publishes to NATS subject: webhook.inbound", WHITE),
    ("Pull worker subscribes with queue group", WHITE),
    ("Worker deserializes the envelope", WHITE),
    ("Router runs 6-layer classification", ORANGE),
    ("Matched agent receives the payload", GREEN),
    ("Agent processes and writes results back to API", GREEN),
]
for j, (text, color) in enumerate(steps):
    add_text(slide, Inches(1.4), Inches(3.7) + j * Inches(0.4),
             Inches(4.7), Inches(0.4),
             f"{j+1}.  {text}", size=16, color=color)

# Right: why NATS
card(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.8))
add_text(slide, Inches(7.4), Inches(3.0), Inches(4.7), Inches(0.5),
         "Why NATS JetStream?", size=22, color=GREEN, bold=True)

reasons = [
    ("Durable delivery", "Messages survive pod restarts"),
    ("Pull-based", "Workers control their own pace"),
    ("Queue groups", "Multiple workers, no duplicate processing"),
    ("Backpressure", "Slow consumers don\u2019t lose messages"),
    ("At-least-once", "Failed processing gets retried"),
]
for j, (title, desc) in enumerate(reasons):
    y = Inches(3.7) + j * Inches(0.55)
    add_text(slide, Inches(7.4), y, Inches(4.7), Inches(0.3),
             title, size=17, color=ACCENT, bold=True)
    add_text(slide, Inches(7.4), y + Inches(0.28), Inches(4.7), Inches(0.3),
             desc, size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 6-Layer Classification
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "6-Layer Classification")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "The router tries each layer in order. First match wins. "
         "Most data is classified without calling an LLM at all.",
         size=19, color=LIGHT_GRAY)

make_table(slide, Inches(1), Inches(2.7), Inches(11.333), [
    ["Layer", "Method", "Example"],
    ["0", "Channel message shortcut",
     "WhatsApp message \u2192 channel_message agent"],
    ["0a", "UDE source_type mapping",
     "source_type=email \u2192 email agent"],
    ["1", "Explicit data_type field",
     "Caller specifies data_type=pdf directly"],
    ["2", "Payload field heuristics (50+ patterns)",
     "Has latitude + longitude \u2192 gps agent"],
    ["3", "MIME registry lookup",
     "application/pdf \u2192 pdf agent"],
    ["4", "URL extension sniffing",
     ".csv \u2192 csv agent, .mp4 \u2192 video_url agent"],
    ["5", "Fallback",
     "Unrecognized content \u2192 binary_blob agent"],
], font_size=15)

add_text(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
         "Layers 0\u20134 are deterministic \u2014 no LLM cost, sub-millisecond. "
         "Only truly ambiguous content falls through to the blob agent.",
         size=17, color=GREEN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 42 Agent Categories
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "42 Specialized Agents")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "Each agent knows how to analyze its data type. The right model tier "
         "is selected automatically \u2014 simple data gets a small model, complex data gets a large one.",
         size=18, color=LIGHT_GRAY)

make_table(slide, Inches(1), Inches(2.7), Inches(5.3), [
    ["Category", "Agents"],
    ["Documents (5)", "PDF, Office, HTML, Markdown, Web Page"],
    ["Media (6)", "Image, Image Batch, Audio/Video URL/Stream"],
    ["Communication (6)", "Chat, Channel Msg, Email, Feed, Calendar, Conferencing"],
    ["Structured (4)", "JSON, CSV, YAML, XML"],
    ["Code & Logs (3)", "Code, Log File, Log Stream"],
    ["Sensor (4)", "GPS, Biometric, IoT Sensor, Sensor"],
], font_size=14)

make_table(slide, Inches(7), Inches(2.7), Inches(5.333), [
    ["Category", "Agents"],
    ["Spatial (3)", "Geospatial, LiDAR, 3D Model"],
    ["Specialized (5)", "Financial, Industrial, Infra, Scientific, Satellite"],
    ["Binary (3)", "Archive, Time Series, Binary Blob"],
    ["Other (3)", "Medical Imaging, Vehicle Telemetry, URL Download"],
], font_size=14)

# Model tiers
card(slide, Inches(7), Inches(5.2), Inches(5.333), Inches(1.6))
add_text(slide, Inches(7.4), Inches(5.3), Inches(4.5), Inches(0.4),
         "5-Tier Smart Routing", size=18, color=ACCENT, bold=True)
tiers = "Small (4b) \u2192 Medium (12b) \u2192 Large (27b) \u2192 Multimodal \u2192 Omni"
add_text(slide, Inches(7.4), Inches(5.8), Inches(4.5), Inches(0.4),
         tiers, size=15, color=LIGHT_GRAY)
add_text(slide, Inches(7.4), Inches(6.2), Inches(4.5), Inches(0.4),
         "80% of data runs on the smallest model", size=15, color=GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What an Agent Produces
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "What AI Enrichment Produces")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "Each agent writes back multiple artifacts. Together, they make "
         "the original data deeply searchable and queryable.",
         size=19, color=LIGHT_GRAY)

artifacts = [
    ("Viewpoint", ACCENT,
     "The agent\u2019s AI analysis of the content \u2014 "
     "a structured summary with input_content, "
     "output_content, and an AISignature tracking "
     "which model produced it and when."),
    ("Embeddings", GREEN,
     "Vector representations for semantic search. "
     "Content is chunked into 1000-char windows "
     "with 200-char overlap. Each chunk gets its "
     "own embedding vector (768+ dimensions)."),
    ("Entities", ORANGE,
     "Extracted people, organizations, locations, "
     "products, dates, URLs, concepts, and events. "
     "Dual-written to Postgres and Graphiti/Neo4j "
     "with temporal valid_at/invalid_at timestamps."),
    ("Relationships", ACCENT2,
     "Directed connections between entities \u2014 "
     "works_at, located_in, mentioned_in, etc. "
     "Stored in the knowledge graph for BFS "
     "traversal during graph search."),
]

for i, (title, color, desc) in enumerate(artifacts):
    left = Inches(1) + (i % 2) * Inches(5.8)
    top = Inches(2.7) + (i // 2) * Inches(2.2)
    card(slide, left, top, Inches(5.4), Inches(1.9), border_color=color)
    add_text(slide, left + Inches(0.4), top + Inches(0.2), Inches(4.6), Inches(0.4),
             title, size=22, color=color, bold=True)
    add_text(slide, left + Inches(0.4), top + Inches(0.7), Inches(4.6), Inches(1.0),
             desc, size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Embedding Pipeline
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "Embedding Pipeline")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "Embeddings power semantic search. Every piece of content is chunked, "
         "vectorized, and stored for cosine similarity retrieval.",
         size=19, color=LIGHT_GRAY)

# Left: chunking
card(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(2.0))
add_text(slide, Inches(1.4), Inches(3.0), Inches(4.7), Inches(0.5),
         "Chunking Strategy", size=22, color=ACCENT, bold=True)
add_text(slide, Inches(1.4), Inches(3.6), Inches(4.7), Inches(1.0),
         "1000-character windows with 200-char overlap.\n"
         "Overlap ensures no information is lost at\n"
         "chunk boundaries. Each chunk is embedded\n"
         "independently.",
         size=16, color=LIGHT_GRAY)

# Right: fallback chain
card(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(2.0))
add_text(slide, Inches(7.4), Inches(3.0), Inches(4.7), Inches(0.5),
         "Embedding Engine Cascade", size=22, color=GREEN, bold=True)

engines = [
    ("1.  Local Ollama", "In-cluster, free, fastest"),
    ("2.  Ollama Cloud", "Fallback if local is down"),
    ("3.  Gemini API", "Final fallback, cloud-based"),
]
for j, (engine, note) in enumerate(engines):
    y = Inches(3.7) + j * Inches(0.35)
    add_text(slide, Inches(7.4), y, Inches(2.5), Inches(0.35),
             engine, size=16, color=ACCENT, bold=True)
    add_text(slide, Inches(9.9), y, Inches(2.5), Inches(0.35),
             note, size=14, color=LIGHT_GRAY)

# Storage
card(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(1.5))
add_text(slide, Inches(1.4), Inches(5.4), Inches(4), Inches(0.5),
         "Where Embeddings Live", size=22, color=ACCENT, bold=True)
add_text(slide, Inches(1.4), Inches(5.9), Inches(10), Inches(0.6),
         "Vectors stored in Postgres via pgvector extension (cosine similarity index). "
         "Each embedding links back to its source data_id and chunk_text, "
         "along with an AISignature recording which model and engine produced it.",
         size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Knowledge Graph
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "Dual-Layer Knowledge Graph")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "Entities and relationships are dual-written to two complementary stores \u2014 "
         "always-on Postgres and optional temporal Graphiti/Neo4j.",
         size=19, color=LIGHT_GRAY)

# Left: Postgres
card(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(3.8), border_color=ACCENT)
add_text(slide, Inches(1.4), Inches(3.0), Inches(4.7), Inches(0.5),
         "Layer 1: Postgres", size=24, color=ACCENT, bold=True)
add_text(slide, Inches(1.4), Inches(3.6), Inches(4.7), Inches(0.3),
         "Always active. Zero extra infrastructure.", size=15, color=WHITE)

items_pg = [
    "8 entity types: person, organization,\nproduct, location, date, URL, concept, event",
    "Directed relationships between entities",
    "Entity \u2192 source data mappings",
    "Deduplicated via canonical form index",
]
for j, item in enumerate(items_pg):
    add_text(slide, Inches(1.4), Inches(4.2) + j * Inches(0.6),
             Inches(4.7), Inches(0.6),
             f"\u2022  {item}", size=15, color=LIGHT_GRAY)

# Right: Graphiti
card(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.8), border_color=GREEN)
add_text(slide, Inches(7.4), Inches(3.0), Inches(4.7), Inches(0.5),
         "Layer 2: Graphiti + Neo4j", size=24, color=GREEN, bold=True)
add_text(slide, Inches(7.4), Inches(3.6), Inches(4.7), Inches(0.3),
         "Optional. Adds temporal reasoning.", size=15, color=WHITE)

items_neo = [
    "Temporal facts with valid_at / invalid_at\ntimestamps for point-in-time queries",
    "LLM-powered entity resolution merges\nduplicates (\"J. Smith\" = \"John Smith\")",
    "Community detection clusters related entities",
    "Fire-and-forget async dual-write from API",
]
for j, item in enumerate(items_neo):
    add_text(slide, Inches(7.4), Inches(4.2) + j * Inches(0.6),
             Inches(4.7), Inches(0.6),
             f"\u2022  {item}", size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 5 Search Modes
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "5 Search Modes")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "All the enrichment feeds into a multi-signal search engine. "
         "Each mode taps different indexes for different strengths.",
         size=19, color=LIGHT_GRAY)

make_table(slide, Inches(1), Inches(2.7), Inches(11.333), [
    ["Mode", "Engine", "Best For"],
    ["vector", "pgvector cosine similarity",
     "Semantic meaning \u2014 \"things like this\""],
    ["fts", "Postgres tsvector / BM25",
     "Exact terms, names, IDs, keywords"],
    ["hybrid", "vector + BM25 via RRF",
     "General purpose (default) \u2014 best of both"],
    ["graph", "Graphiti BFS + semantic on Neo4j",
     "Entity relationships, temporal facts"],
    ["full", "hybrid + graph merged with RRF",
     "Maximum recall across all signals"],
], font_size=16)

add_text(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(0.6),
         "Hybrid merges vector and keyword scores via Reciprocal Rank Fusion. "
         "Full mode adds graph results on top for the broadest coverage.",
         size=17, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Reranking
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "4 Reranking Strategies")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "After initial retrieval, results can be re-scored to improve "
         "relevance, diversity, or precision before reaching the user.",
         size=19, color=LIGHT_GRAY)

rerankers = [
    ("None", LIGHT_GRAY,
     "Raw scores from the search engine.\nFastest option \u2014 good for simple queries."),
    ("RRF", ACCENT,
     "Reciprocal Rank Fusion. Merges ranked\nlists from multiple search modes fairly."),
    ("MMR", ORANGE,
     "Maximal Marginal Relevance. Prioritizes\ndiversity \u2014 avoids redundant results."),
    ("Cross-encoder", GREEN,
     "LLM re-scores each result against the query.\nHighest accuracy, highest cost."),
]

for i, (title, color, desc) in enumerate(rerankers):
    left = Inches(1) + i * Inches(3.0)
    card(slide, left, Inches(2.8), Inches(2.7), Inches(2.5), border_color=color)
    add_text(slide, left + Inches(0.3), Inches(3.0), Inches(2.1), Inches(0.5),
             title, size=24, color=color, bold=True)
    add_text(slide, left + Inches(0.3), Inches(3.6), Inches(2.1), Inches(1.4),
             desc, size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RAG Chat
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "RAG Chat with Citations")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "The final step: turn search results into a conversational answer "
         "with traceable inline citations back to source data.",
         size=19, color=LIGHT_GRAY)

card(slide, Inches(1), Inches(2.8), Inches(11.333), Inches(3.8))
add_text(slide, Inches(1.4), Inches(3.0), Inches(10), Inches(0.5),
         "POST /api/v1/ai/query/chat", size=18, color=GREEN, font="Courier New")

steps = [
    ("1", "User asks a question", ACCENT),
    ("2", "Embed the query using the same embedding engine", ACCENT),
    ("3", "Run selected search mode (vector, hybrid, graph, full)", ACCENT),
    ("4", "Apply reranker to retrieved results", ACCENT),
    ("5", "Build numbered context:  [1] source A,  [2] source B,  [3] source C", ORANGE),
    ("6", "Send context + conversation history + question to LLM", ORANGE),
    ("7", "LLM generates answer with inline [1] [2] [3] citation markers", GREEN),
    ("8", "Extract cited indices, return only referenced sources", GREEN),
]

for j, (num, text, color) in enumerate(steps):
    y = Inches(3.6) + j * Inches(0.38)
    add_text(slide, Inches(1.8), y, Inches(0.4), Inches(0.35),
             num, size=17, color=color, bold=True)
    add_text(slide, Inches(2.3), y, Inches(9.5), Inches(0.35),
             text, size=17, color=WHITE if color != GREEN else GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — End-to-End Example
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "End-to-End Example")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
         "A Slack message about a product launch, from ingestion to query answer.",
         size=19, color=LIGHT_GRAY)

# Left column: ingestion
card(slide, Inches(1), Inches(2.6), Inches(5.5), Inches(4.2))
add_text(slide, Inches(1.4), Inches(2.8), Inches(4.7), Inches(0.5),
         "Ingestion + Enrichment", size=22, color=ACCENT, bold=True)

ing_steps = [
    ("Slack adapter normalizes to UniversalEnvelope", WHITE),
    ("API stores raw message + publishes to NATS", WHITE),
    ("Router: channel message \u2192 layer 0 match", ORANGE),
    ("ChannelMessageAgent analyzes with Gemma3:4b", WHITE),
    ("Viewpoint: intent=announcement, topic=launch", GREEN),
    ("Embeddings: 2 chunks \u00d7 768-dim vectors", GREEN),
    ("Entities: \"Product X\", \"Q3\", \"pricing model\"", GREEN),
    ("Graph: Product X \u2192 launches_in \u2192 Q3", GREEN),
]
for j, (text, color) in enumerate(ing_steps):
    add_text(slide, Inches(1.4), Inches(3.5) + j * Inches(0.38),
             Inches(4.7), Inches(0.38),
             f"\u2022  {text}", size=14, color=color)

# Right column: query
card(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(4.2))
add_text(slide, Inches(7.4), Inches(2.8), Inches(4.7), Inches(0.5),
         "Query + Answer", size=22, color=GREEN, bold=True)

query_steps = [
    ("User asks: \"When is the product launch?\"", WHITE),
    ("Query embedded \u2192 hybrid search (vector + FTS)", WHITE),
    ("Vector match: Slack msg (0.89 similarity)", ORANGE),
    ("FTS match: \"product launch\" keyword hit", ORANGE),
    ("RRF merges scores \u2192 Slack msg ranked #1", ORANGE),
    ("Graph boost: Product X entity context added", ORANGE),
    ("LLM: \"Based on [1], the product launch is", GREEN),
    ("scheduled for Q3. Action items include...\"", GREEN),
]
for j, (text, color) in enumerate(query_steps):
    add_text(slide, Inches(7.4), Inches(3.5) + j * Inches(0.38),
             Inches(4.7), Inches(0.38),
             f"\u2022  {text}", size=14, color=color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Summary
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()

add_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.8),
         "From Ingestion to Intelligence", size=48, bold=True,
         align=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(5), Inches(1.9), Inches(3.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

capabilities = [
    ("Ingest", "from 33 channel adapters, direct API, or universal envelope"),
    ("Classify", "with 6 deterministic layers \u2014 no LLM needed for most data"),
    ("Route", "to 1 of 42 specialized agents matched to the data type"),
    ("Analyze", "with the smallest model that can handle the complexity"),
    ("Embed", "into 768+ dim vectors with cascading engine fallbacks"),
    ("Extract", "entities and relationships into a temporal knowledge graph"),
    ("Search", "across 5 modes \u2014 vector, keyword, hybrid, graph, full"),
    ("Answer", "with RAG chat, conversation history, and [1][2] citations"),
]

for j, (verb, rest) in enumerate(capabilities):
    top = Inches(2.4) + j * Inches(0.5)
    tf = add_text(slide, Inches(2.5), top, Inches(8.333), Inches(0.45),
                  "", size=20, align=PP_ALIGN.LEFT)
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{verb}  "
    run1.font.size = Pt(20)
    run1.font.color.rgb = ACCENT
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = rest
    run2.font.size = Pt(20)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = "Calibri"

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.6),
         "Every piece of data, automatically enriched and instantly queryable.",
         size=26, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
output_path = "/Users/paragagarwal/repos/mem-dog/docs/ingestion-pipeline.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
