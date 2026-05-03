#!/usr/bin/env python3
"""Generate memdog presentation as .pptx from presentation.md content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy background
ACCENT = RGBColor(0x38, 0xBD, 0xF8)         # Bright blue accent
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)        # Purple accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
TABLE_HEADER_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_ROW_BG = RGBColor(0x15, 0x1F, 0x32)
TABLE_ALT_BG = RGBColor(0x1A, 0x25, 0x38)
GREEN = RGBColor(0x34, 0xD3, 0x99)
ORANGE = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        # Add bullet character
        p.text = f"\u2022  {item}"
        # Color the bullet
    return tf


def make_table(slide, left, top, width, rows_data, col_widths=None, font_size=13):
    """Create a styled table. rows_data[0] is header row."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    table_height = Inches(0.4 * num_rows)
    shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)

            # Style
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = ACCENT
                else:
                    paragraph.font.color.rgb = WHITE

            # Background
            fill = cell.fill
            fill.solid()
            if r_idx == 0:
                fill.fore_color.rgb = TABLE_HEADER_BG
            elif r_idx % 2 == 0:
                fill.fore_color.rgb = TABLE_ALT_BG
            else:
                fill.fore_color.rgb = TABLE_ROW_BG

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "Mem-Dog", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(3.1), Inches(3.333))

add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
             "The Private AI System", font_size=32, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
             "AI should be private. Not as a feature \u2014 as the foundation.",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
             "42 AI Agents  \u00b7  33 Channel Adapters  \u00b7  300+ App Integrations  \u00b7  5 Search Modes  \u00b7  $0/month",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "The Problem", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

# Three columns
col_w = Inches(3.6)
col_gap = Inches(0.4)
col_start = Inches(0.8)

for i, (title, body) in enumerate([
    ("For Individuals",
     "Personal thoughts, health records, financial details \u2014 one breach away from exposure. You pay per token, per seat, per month."),
    ("For Enterprises",
     "Trade secrets and customer data flow through third-party systems with opaque retention policies. Compliance teams can\u2019t audit what they don\u2019t control."),
    ("For Developers",
     "Memory SDK here, integration platform there, vector DB somewhere else, knowledge graph from another vendor \u2014 none of them talk to each other.")
]):
    left = col_start + i * (col_w + col_gap)
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), col_w, Inches(4.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TABLE_HEADER_BG
    shape.line.fill.background()

    add_text_box(slide, left + Inches(0.3), Inches(1.9), col_w - Inches(0.6), Inches(0.5),
                 title, font_size=22, color=ACCENT, bold=True)
    add_text_box(slide, left + Inches(0.3), Inches(2.6), col_w - Inches(0.6), Inches(2.8),
                 body, font_size=16, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
             "We built Mem-Dog to end that trade-off.",
             font_size=24, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — What is Mem-Dog?
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "What is Mem-Dog?", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.0),
             "A complete, self-hosted AI system that ingests data from 33 channel adapters and 300+ apps, "
             "enriches it with 42 specialized AI agents, stores it in a temporal knowledge graph, and makes "
             "it searchable through 5 search modes with RAG chat \u2014 all running on hardware you control.",
             font_size=17, color=LIGHT_GRAY)

# Four pillars
pillars = [
    ("\U0001f512 Private by Design",
     "Data never leaves your network. No cloud dependency. Full offline/air-gapped capability."),
    ("\u26a1 Fast Locally",
     "6-layer deterministic classification. Most queries in milliseconds. 80% runs on the 4b model."),
    ("\U0001f4b0 Cost Efficient",
     "Local Ollama models. $0/month on a Mac Mini. 60-80% cheaper than cloud-only at scale."),
    ("\U0001f9e0 Genuinely Smart",
     "42 agents, dual-layer knowledge graph, 5 search modes, 4 rerankers, RAG chat with citations.")
]

for i, (title, desc) in enumerate(pillars):
    left = Inches(0.8) + i * Inches(3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.8), Inches(2.8), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TABLE_HEADER_BG
    shape.line.fill.background()

    add_text_box(slide, left + Inches(0.2), Inches(3.1), Inches(2.4), Inches(0.7),
                 title, font_size=18, color=ACCENT, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(3.8), Inches(2.4), Inches(2.2),
                 desc, font_size=14, color=LIGHT_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — The Complete Platform
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "The Complete Platform", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

make_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), [
    ["Layer", "What Mem-Dog Provides", "What Others Offer"],
    ["Ingestion", "33 channel adapters + 300+ app integrations (Nango)", "Mem0: API-only. Zep: SDK-only."],
    ["Enrichment", "42 typed AI agents, 6-layer classification, 60+ data types", "Mem0: fact extraction only."],
    ["Storage", "10 memory types, versioned mutations, 3 backends, per-item ACLs", "Mem0: 3 types. Zep: messages only."],
    ["Knowledge Graph", "Dual-layer Postgres + Graphiti/Neo4j with temporal reasoning", "Zep: Graphiti only. Mem0: none."],
    ["Search", "5 modes (vector, FTS, hybrid, graph, full) + 4 rerankers", "Mem0: vector only. Zep: semantic + graph."],
    ["Query", "RAG chat with inline [1][2] citations, memory-scoped", "Dify: chat builder. Others: BYO."],
    ["Agent", "DigiMe conversational AI in 25+ messaging apps", "No equivalent in any competitor."],
    ["AI Studio", "Search, Models, Routing, Agents, Infrastructure \u2014 all in one UI", "No equivalent."],
    ["UI", "Full web platform \u2014 dashboard, playground, telemetry, settings", "Mem0: none. Zep: none."],
    ["SDKs", "Python, TypeScript, Go, Rust, Ruby + LangChain/CrewAI/MCP", "Varies."],
], font_size=12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Architecture", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Architecture boxes - simplified flow diagram
components = [
    (Inches(0.5), Inches(1.6), Inches(2.5), "300+ Apps (Nango)\n33 Channel Adapters\nSlack, WhatsApp, Email...", ACCENT),
    (Inches(0.5), Inches(3.2), Inches(2.5), "Webhook Gateway\nNormalize \u2192 Universal\nEnvelope Format", ACCENT2),
    (Inches(3.8), Inches(1.6), Inches(2.5), "Web UI\nNext.js 14 / React 18\nAI Studio", GREEN),
    (Inches(3.8), Inches(3.2), Inches(2.5), "API (FastAPI)\n33 routers, 80+ endpoints\nPython 3.12", GREEN),
    (Inches(7.1), Inches(1.6), Inches(2.8), "Webhook Pipeline\nNATS JetStream\n42 AI Agents\n6-Layer Classification", ORANGE),
    (Inches(7.1), Inches(3.2), Inches(2.8), "Ollama (3 tiers)\nGemma3 4b/12b/27b\nQwen3-VL, Qwen3.5", ORANGE),
    (Inches(10.5), Inches(1.6), Inches(2.3), "Supabase\nPostgres 16 + pgvector\nGoTrue Auth", ACCENT),
    (Inches(10.5), Inches(3.2), Inches(2.3), "Neo4j + Graphiti\nTemporal Knowledge\nGraph", ACCENT),
]

for left, top, width, text, accent_color in components:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TABLE_HEADER_BG
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(1.5)

    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(text.split("\n")):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = "Calibri"
        p.font.color.rgb = WHITE if j == 0 else LIGHT_GRAY
        p.font.bold = (j == 0)
        p.alignment = PP_ALIGN.CENTER

# Component table
make_table(slide, Inches(0.5), Inches(5.0), Inches(12.3), [
    ["Component", "Stack", "Port", "Deployment"],
    ["API", "Python 3.12, FastAPI, 33 routers", "8080", "GKE (memdog)"],
    ["UI", "Next.js 14, React 18, TypeScript", "3000", "Cloud Run"],
    ["Pipeline", "NATS JetStream, ADK, 42 agents", "8080", "GKE (webhook-pipeline)"],
    ["Gateway", "FastAPI, LiteLLM, 33 adapters", "8080", "GKE (webhook-gateway)"],
    ["DigiMe", "Node.js, OpenClaw, 4 skills", "18789", "GKE (webhook-gateway)"],
], font_size=11)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Multi-Channel Ingestion
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Multi-Channel Ingestion", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Three ingestion paths
paths = [
    ("Path 1: Channel Ingestion",
     "33 channel adapters",
     ["Channel delivers to Webhook Gateway", "Normalize \u2192 UniversalEnvelope",
      "Forward to API \u2192 Store + NATS", "6-layer classify \u2192 Route to agent",
      "LLM analysis \u2192 Viewpoint + Embedding"]),
    ("Path 2: Direct Ingestion",
     "UI or REST API",
     ["POST /api/v1/data with metadata", "Write to Graphiti as episode",
      "Optional: forward to pipeline", "Supports 6 upload modes in UI", ""]),
    ("Path 3: Conversational",
     "DigiMe Agent (25+ apps)",
     ["User messages DigiMe on any platform", "memdog-bridge forwards to pipeline",
      "memdog-ingest stores explicit data", "Channel identity resolves user_id", ""]),
]

for i, (title, subtitle, steps) in enumerate(paths):
    left = Inches(0.5) + i * Inches(4.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(3.9), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TABLE_HEADER_BG
    shape.line.fill.background()

    add_text_box(slide, left + Inches(0.3), Inches(1.6), Inches(3.3), Inches(0.5),
                 title, font_size=18, color=ACCENT, bold=True)
    add_text_box(slide, left + Inches(0.3), Inches(2.1), Inches(3.3), Inches(0.4),
                 subtitle, font_size=14, color=LIGHT_GRAY)

    for j, step in enumerate(steps):
        if step:
            add_text_box(slide, left + Inches(0.3), Inches(2.7) + j * Inches(0.45), Inches(3.3), Inches(0.4),
                         f"\u2192  {step}", font_size=13, color=WHITE)

# Channel list at bottom
add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
             "33 Native Adapters: Slack, WhatsApp, Telegram, Email, Discord, GitHub, Jira, Linear, "
             "Salesforce, HubSpot, Zoom, MS Teams, Notion, Stripe, PagerDuty, Sentry, Datadog, "
             "and 16 more  |  15+ Bridge Channels via OpenClaw: Signal, Matrix, IRC, Google Chat, Line...",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — 42-Agent AI Pipeline
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "42-Agent AI Enrichment Pipeline", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Processing flow
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
             "Ingest \u2192 Stage \u2192 Classify (6-layer) \u2192 Route \u2192 Analyze (LLM) \u2192 Viewpoint \u2192 Embedding \u2192 Entity Extraction \u2192 Graph",
             font_size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# 6-layer classification
make_table(slide, Inches(0.5), Inches(2.0), Inches(6.0), [
    ["Layer", "Method", "Example"],
    ["0", "URL not-downloaded detection", "URL \u2192 url_download agent first"],
    ["1", "Channel message detection", "WhatsApp msg \u2192 channel_message"],
    ["2", "Explicit data_type field", "Direct classification from caller"],
    ["3", "Payload field heuristic", "Has lat/lng \u2192 gps"],
    ["4", "MIME registry", "application/pdf \u2192 pdf"],
    ["5", "URL extension sniff", ".csv \u2192 csv"],
    ["6", "LLM classifier (fallback)", "Gemma3:4b classifies ambiguous"],
], font_size=11)

# Agent categories
make_table(slide, Inches(6.8), Inches(2.0), Inches(6.0), [
    ["Category", "Count", "Agents"],
    ["Documents", "5", "PDF, Office, HTML, Markdown, Web Page"],
    ["Media", "6", "Image, Image Batch, Audio/Video URL/Stream"],
    ["Communication", "6", "Chat, Channel Msg, Email, Feed, Calendar, Conferencing"],
    ["Structured", "4", "JSON, CSV, YAML, XML"],
    ["Code & Logs", "3", "Code, Log File, Log Stream"],
    ["Sensor & IoT", "4", "GPS, Biometric, IoT Sensor, Sensor"],
    ["Spatial", "3", "Geospatial, LiDAR, 3D Model"],
    ["Specialized", "5", "Financial, Industrial, Infra, Scientific, Satellite"],
    ["Binary", "3", "Binary Blob, Archive, Time Series"],
    ["Other", "3", "Medical Imaging, Vehicle Telemetry, URL Download"],
], font_size=11)

add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
             "Most data classified deterministically in layers 0-5 \u2014 avoiding LLM calls entirely (saves cost + latency)",
             font_size=14, color=GREEN, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — AI Studio
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "AI Studio", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
             "Unified control center for all AI capabilities", font_size=18, color=LIGHT_GRAY)

make_table(slide, Inches(0.5), Inches(1.9), Inches(12.3), [
    ["Tab", "What It Does"],
    ["Search & Chat", "RAG search with 5 modes + 4 rerankers. Viewpoint browser. Embedding manager. Knowledge chat with [1][2] citations."],
    ["Models", "Model Garden \u2014 connect 16+ AI providers. Encrypted API key storage. Model catalog. Per-user provider config."],
    ["Routing", "Smart Routing \u2014 per-agent-type model assignment. Primary + fallback chains. Visual routing table."],
    ["Agents", "Processing Flags (enable/disable per type) + Agent Configs (custom prompts, schemas, model overrides)."],
    ["Infrastructure", "K8s Pod Management \u2014 create, scale (0-3), delete Ollama pods from browser. Logs, metrics, status."],
], font_size=13)

# Model tiers
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
             "5-Tier Smart Routing", font_size=22, color=WHITE, bold=True)

make_table(slide, Inches(0.5), Inches(5.1), Inches(6.0), [
    ["Tier", "Default Model", "Used For"],
    ["Small", "Gemma3:4b", "JSON, CSV, YAML, IoT, simple text"],
    ["Medium", "Gemma3:12b", "Code, email, chat, financial"],
    ["Large", "Gemma3:27b", "PDFs, Office docs, complex documents"],
    ["Multimodal", "Qwen3-VL", "Images, visual PDFs, screenshots"],
    ["Omni", "Qwen3.5", "Audio, video, multi-format"],
], font_size=12)

# Providers
add_text_box(slide, Inches(7.0), Inches(4.5), Inches(5), Inches(0.5),
             "16+ AI Providers", font_size=22, color=WHITE, bold=True)

add_text_box(slide, Inches(7.0), Inches(5.1), Inches(5.5), Inches(2.0),
             "Ollama (local) \u00b7 Ollama Cloud \u00b7 Google Gemini \u00b7 OpenAI \u00b7 Anthropic \u00b7 "
             "Groq \u00b7 Mistral \u00b7 Cohere \u00b7 DeepSeek \u00b7 xAI \u00b7 OpenRouter \u00b7 "
             "Together AI \u00b7 HuggingFace \u00b7 AWS Bedrock \u00b7 vLLM \u00b7 LiteLLM",
             font_size=14, color=LIGHT_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Knowledge Graph
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Temporal Knowledge Graph", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Two layers side by side
for i, (title, subtitle, items) in enumerate([
    ("Layer 1: Postgres", "Always Active, Zero Infrastructure",
     ["8 entity types: person, org, product, location, date, URL, concept, event",
      "Directed relationships (works_at, located_in, etc.)",
      "Entity \u2192 source data mappings",
      "Deduplicated via canonical form unique index",
      "No additional infra \u2014 same Postgres"]),
    ("Layer 2: Graphiti + Neo4j", "Optional, Temporal Reasoning",
     ["Temporal facts with valid_at / invalid_at timestamps",
      "Point-in-time queries (\"Who was CEO in 2024?\")",
      "LLM-powered entity resolution & dedup",
      "Community detection via label propagation",
      "Fire-and-forget dual-write from API"]),
]):
    left = Inches(0.5) + i * Inches(6.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(6.1), Inches(4.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TABLE_HEADER_BG
    shape.line.fill.background()

    add_text_box(slide, left + Inches(0.3), Inches(1.6), Inches(5.5), Inches(0.5),
                 title, font_size=22, color=ACCENT, bold=True)
    add_text_box(slide, left + Inches(0.3), Inches(2.1), Inches(5.5), Inches(0.4),
                 subtitle, font_size=14, color=LIGHT_GRAY)

    for j, item in enumerate(items):
        add_text_box(slide, left + Inches(0.3), Inches(2.7) + j * Inches(0.42), Inches(5.5), Inches(0.4),
                     f"\u2022  {item}", font_size=13, color=WHITE)

# Entity-aware RAG
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.2),
             "Entity-Aware RAG: Query \u2192 Match entities \u2192 BFS graph traversal \u2192 "
             "Inject entity context into LLM \u2192 Return enriched answers with structured knowledge",
             font_size=15, color=GREEN, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Search & RAG
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "5 Search Modes + 4 Rerankers", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

make_table(slide, Inches(0.5), Inches(1.4), Inches(12.3), [
    ["Mode", "Engine", "How It Works", "Best For"],
    ["vector", "pgvector", "Cosine similarity over embeddings", "Semantic meaning"],
    ["fts", "Postgres tsvector", "BM25 keyword matching", "Exact terms, names, IDs"],
    ["hybrid", "pgvector + tsvector", "Vector + BM25 via Reciprocal Rank Fusion", "General purpose (default)"],
    ["graph", "Graphiti + Neo4j", "BFS traversal + semantic search on knowledge graph", "Entity relationships, temporal"],
    ["full", "All engines", "Hybrid + Graph in parallel, RRF merged", "Maximum recall"],
], font_size=13)

make_table(slide, Inches(0.5), Inches(4.0), Inches(12.3), [
    ["Reranker", "Method", "When To Use"],
    ["None", "Raw scores from search engine", "Fast, good for simple queries"],
    ["RRF", "Reciprocal Rank Fusion", "Merging results from multiple modes"],
    ["MMR", "Maximal Marginal Relevance", "Diverse results, not just top-N similar"],
    ["Cross-encoder", "LLM-scored relevance", "Highest accuracy, model re-scores each result"],
], font_size=13)

# RAG chat box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.3))
shape.fill.solid()
shape.fill.fore_color.rgb = TABLE_HEADER_BG
shape.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(5.9), Inches(3), Inches(0.5),
             "RAG Chat with Citations", font_size=20, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
             "Search \u2192 Rerank \u2192 Build numbered context \u2192 LLM + conversation history \u2192 "
             "Answer with inline [1] [2] [3] citation markers \u2192 Memory-scoped results",
             font_size=14, color=LIGHT_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Memory System
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Memory System", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

make_table(slide, Inches(0.5), Inches(1.3), Inches(7.0), [
    ["Category", "Type", "Default TTL"],
    ["Conversation", "conversation", "1 hour"],
    ["Session", "session", "24 hours"],
    ["Session", "timeline", "7 days"],
    ["Session", "tracing", "3 days"],
    ["User", "user", "Never"],
    ["User", "factual", "Never"],
    ["User", "episodic", "Temporary"],
    ["User", "semantic", "Never"],
    ["User", "custom", "Varies"],
    ["Organizational", "organizational", "Never"],
], font_size=12)

# Features
features = [
    ("Versioning", "Every mutation creates a new version with diff tracking"),
    ("Access Control", "4 levels: private, shared, public, restricted + shared_with list"),
    ("TTL & Expiry", "Default TTL per type. Override with ttl_hours or no_expiry=true"),
    ("LLM Compression", "Summarize verbose memories, archive originals, keep key facts"),
    ("Project Scoping", "Memories linked to org/project hierarchy for multi-tenant isolation"),
    ("ID Format", "mem_<type>_<ulid> \u2014 e.g., mem_conv_01JQ, mem_user_01JQ"),
]

add_text_box(slide, Inches(8.0), Inches(1.3), Inches(4.8), Inches(0.5),
             "Features", font_size=22, color=ACCENT, bold=True)

for j, (feat, desc) in enumerate(features):
    add_text_box(slide, Inches(8.0), Inches(1.9) + j * Inches(0.7), Inches(4.8), Inches(0.7),
                 f"{feat}", font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.0), Inches(2.2) + j * Inches(0.7), Inches(4.8), Inches(0.4),
                 desc, font_size=12, color=LIGHT_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — DigiMe
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "DigiMe: Your Personal AI Assistant", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
             "An AI agent that lives in your messaging apps \u2014 turning any chat into a window to your knowledge base.",
             font_size=17, color=LIGHT_GRAY)

# Skills table
make_table(slide, Inches(0.5), Inches(2.0), Inches(12.3), [
    ["Skill", "Purpose", "Behavior"],
    ["memdog-bridge", "Record everything", "Forwards EVERY message to pipeline. Non-negotiable."],
    ["memdog-ingest", "Store explicit data", "Stores notes, facts, session memories via POST /api/v1/data"],
    ["memdog-query", "Lookup data", "Retrieves memories, data items, lists via API"],
    ["memdog-semantic-search", "Recall by meaning", "Vector similarity search via POST /api/v1/ai/query/semantic"],
], font_size=13)

# Multi-user isolation
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = TABLE_HEADER_BG
shape.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5.2), Inches(0.5),
             "Multi-User Data Isolation", font_size=20, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5.2), Inches(1.5),
             "User A (WhatsApp) \u2192 channel identity \u2192 user_id_A\n"
             "User B (Signal) \u2192 channel identity \u2192 user_id_B\n\n"
             "All queries scoped to user_id. User A cannot see User B's data.",
             font_size=14, color=LIGHT_GRAY)

# Platforms
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(4.0), Inches(6.1), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = TABLE_HEADER_BG
shape.line.fill.background()

add_text_box(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(0.5),
             "25+ Messaging Platforms", font_size=20, color=ACCENT, bold=True)
add_text_box(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(1.5),
             "WhatsApp \u00b7 Telegram \u00b7 Signal \u00b7 Slack \u00b7 Discord \u00b7 "
             "Matrix \u00b7 MS Teams \u00b7 IRC \u00b7 Google Chat \u00b7 Line \u00b7 "
             "Feishu \u00b7 Mattermost \u00b7 Nextcloud Talk \u00b7 Nostr \u00b7 "
             "Twitch \u00b7 Zalo \u00b7 iMessage \u00b7 Rocket.Chat \u00b7 WeChat \u00b7 "
             "Viber \u00b7 Webchat \u00b7 and more",
             font_size=14, color=LIGHT_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Integration Platform
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Integration Platform", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Nango provides vs Mem-Dog adds
make_table(slide, Inches(0.5), Inches(1.3), Inches(5.8), [
    ["Nango Provides", "Details"],
    ["300+ provider templates", "OAuth flows, API schemas, webhooks"],
    ["OAuth2 flows", "Auth code grant, PKCE, state mgmt"],
    ["Auto token refresh", "Transparent, no retry logic needed"],
    ["Credential encryption", "AES-256-GCM"],
    ["Per-user connections", "end_user_id for multi-tenant"],
], font_size=12)

make_table(slide, Inches(6.7), Inches(1.3), Inches(6.1), [
    ["Mem-Dog Adds", "Details"],
    ["Per-user webhooks", "whk_<ulid> with HMAC secrets"],
    ["33 channel adapters", "Platform \u2192 UniversalEnvelope"],
    ["42 AI agents", "Auto-enrich all ingested data"],
    ["Knowledge graph", "Dual-write Postgres + Graphiti"],
    ["Credential proxy", "POST /proxy/{provider}/{path}"],
    ["Response normalize", "?normalize=contact|calendar_event"],
    ["Write-back", "Google Drive, Gmail, Zoom routers"],
], font_size=12)

add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5),
             "Categories: CRM \u00b7 Email \u00b7 Messaging \u00b7 Project Management \u00b7 Cloud Storage \u00b7 "
             "Social Media \u00b7 Developer Tools \u00b7 Analytics \u00b7 E-Commerce \u00b7 Payments \u00b7 "
             "HR \u00b7 Accounting \u00b7 Marketing \u00b7 Support \u00b7 and more",
             font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Web UI
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Web UI", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

make_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), [
    ["Section", "Tabs", "What They Do"],
    ["Knowledge", "Data, Memories", "Browse ingested data. Manage memories with TTL, ACLs, versioning."],
    ["AI Studio", "Search, Models, Routing, Agents, Infra", "Semantic search, model mgmt, smart routing, agent config, K8s pods."],
    ["Monitor", "Insights, Telemetry", "Stats dashboard. OpenTelemetry waterfall trace viewer."],
    ["Develop", "Playground", "Channel Simulator, Data Insert (6 modes), Knowledge Chat, MCP tester."],
    ["Settings", "Profile, Orgs, Apps, Webhooks, Keys", "User profile, org hierarchy, Nango integrations, webhook mgmt, API keys."],
], font_size=13)

# Playground detail
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.5),
             "Interactive Playground", font_size=22, color=WHITE, bold=True)

playground_items = [
    "Channel Simulator \u2014 test webhooks across 33 channels",
    "6-Mode Upload \u2014 text, file, URL, camera, voice, video",
    "Knowledge Chat \u2014 RAG with mode/reranker selectors",
    "MCP Playground \u2014 test 8 MCP tools interactively",
]
add_bullet_list(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5),
                playground_items, font_size=14, color=LIGHT_GRAY)

# Tech stack
add_text_box(slide, Inches(7.0), Inches(4.0), Inches(5), Inches(0.5),
             "Tech Stack", font_size=22, color=WHITE, bold=True)

tech_items = [
    "Next.js 14 + React 18 + TypeScript",
    "Tailwind CSS \u2014 dark-mode glassmorphism",
    "Server-side proxy rewrites to API/Auth",
    "NEXT_PUBLIC_* baked at Docker build time",
]
add_bullet_list(slide, Inches(7.0), Inches(4.6), Inches(5.5), Inches(2.5),
                tech_items, font_size=14, color=LIGHT_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Developer Experience
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Developer Experience", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# SDKs
make_table(slide, Inches(0.5), Inches(1.3), Inches(6.0), [
    ["Language", "Implementation", "Notes"],
    ["Python", "httpx, async", "70+ methods + Simple facade (8 methods)"],
    ["TypeScript", "Native fetch", "Full API coverage, type-safe"],
    ["Go", "stdlib net/http", "Idiomatic, no external deps"],
    ["Rust", "async tokio + reqwest", "Fully async, zero-cost"],
    ["Ruby", "native HTTP", "Clean Ruby interface"],
], font_size=12)

# Agent adapters
make_table(slide, Inches(6.8), Inches(1.3), Inches(6.0), [
    ["Framework", "Adapter"],
    ["LangChain", "ChatMessageHistory + Retriever"],
    ["CrewAI", "save/search persistent memory"],
    ["OpenAI", "Function calling tool"],
    ["MCP", "8 tools for Claude Desktop, Cursor"],
], font_size=12)

# MCP tools
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.5),
             "MCP Server \u2014 8 Tools", font_size=22, color=WHITE, bold=True)

make_table(slide, Inches(0.5), Inches(4.6), Inches(6.0), [
    ["Tool", "What It Does"],
    ["mem_dog_add", "Store new data with tags/metadata"],
    ["mem_dog_search", "Semantic search with mode selection"],
    ["mem_dog_get", "Retrieve a specific data item"],
    ["mem_dog_list", "List data with filters"],
    ["mem_dog_delete", "Delete data"],
    ["mem_dog_entities", "Query knowledge graph entities"],
    ["mem_dog_memories", "List and manage memories"],
    ["mem_dog_chat", "RAG chat with citations"],
], font_size=11)

# Code sample
add_text_box(slide, Inches(7.0), Inches(4.0), Inches(5.5), Inches(0.5),
             "Simple SDK \u2014 8 Methods, Full Power", font_size=18, color=WHITE, bold=True)

code_text = (
    'from mem_dog_client import MemDog\n'
    'm = MemDog(url, api_key="md_...", user_id="u1")\n\n'
    'm.add("Meeting notes", tags=["standup"])\n'
    'results = m.search("standup?", use_ai=True)\n'
    'entities = m.entities("Google")\n'
    'item = m.get("data_01ABC...")\n'
    'm.compress("mem_session_xyz")\n'
    'related = m.related("data_01ABC...")'
)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.6), Inches(6.0), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
shape.line.fill.background()

add_text_box(slide, Inches(7.0), Inches(4.7), Inches(5.5), Inches(2.3),
             code_text, font_size=11, color=GREEN, font_name="Courier New")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Security & Privacy
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Security & Data Privacy", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Security layers
make_table(slide, Inches(0.5), Inches(1.3), Inches(6.0), [
    ["Layer", "Implementation"],
    ["Authentication", "Supabase Auth (email + Google OAuth)"],
    ["API Keys", "Per-user md_* keys, O(1) lookup"],
    ["JWT", "Supabase sub claim, auto-profile"],
    ["Credentials", "AES-256-GCM (Nango) + Fernet (AI keys)"],
    ["Data Isolation", "user_id scoping on every table/query"],
    ["Access Control", "4 levels + shared_with per item"],
    ["Gateway", "API key, IP allowlist, rate limiting"],
    ["Webhooks", "HMAC per endpoint"],
    ["K8s RBAC", "Pod mgmt restricted to namespace"],
], font_size=12)

# Zero trust diagram
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6.0), Inches(4.0))
shape.fill.solid()
shape.fill.fore_color.rgb = TABLE_HEADER_BG
shape.line.fill.background()

add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5.4), Inches(0.5),
             "Zero Trust (Local Mode)", font_size=20, color=GREEN, bold=True)

zero_trust_flow = [
    "User (WhatsApp/Signal/Slack) \u2014 encrypted channel",
    "\u2193  OpenClaw Node (your pod)",
    "\u2193  Webhook Gateway (your pod, identity resolution)",
    "\u2193  42 AI Agents + Ollama (your pod, LOCAL model)",
    "\u2193  Postgres + Neo4j (your pod, your disk)",
    "\u2193  Search/RAG (your pod, local pgvector)",
    "",
    "Data never leaves your infrastructure.",
    "Zero third-party API calls in full local mode.",
]

for j, line in enumerate(zero_trust_flow):
    color = GREEN if "never" in line.lower() or "zero" in line.lower() else LIGHT_GRAY
    add_text_box(slide, Inches(7.1), Inches(2.1) + j * Inches(0.33), Inches(5.4), Inches(0.35),
                 line, font_size=12, color=color)

# Privacy scenarios
make_table(slide, Inches(0.5), Inches(5.6), Inches(12.3), [
    ["Scenario", "How Mem-Dog Helps"],
    ["Healthcare (HIPAA)", "Mac Mini in clinic. Patient data never leaves building. Local Ollama."],
    ["Legal (privilege)", "Attorney-client data processed locally. No third party sees content."],
    ["Corporate IP", "Roadmaps, financials in local Postgres. 42 agents classify without external calls."],
    ["EU (GDPR)", "Host in Germany. Data residency by physics. DELETE = provable erasure."],
    ["Air-gapped", "Full stack runs without internet. Ollama, Postgres, Neo4j \u2014 all local."],
], font_size=12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — Cost Advantage
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Cost Advantage", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

make_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), [
    ["Cost Category", "Mem-Dog (Mac Mini)", "Mem-Dog (GKE)", "Cloud SaaS"],
    ["LLM inference", "$0 (local Ollama)", "$15-50 (Gemini)", "~$160/mo"],
    ["Embeddings", "$0 (local)", "$5-15", "~$20/mo"],
    ["Platform fee", "$0", "$0", "$99+/mo"],
    ["Infrastructure", "$0 (one-time HW)", "~$150-300/mo", "Included"],
    ["Per-seat fee", "$0", "$0", "Per-user pricing"],
    ["Total", "$0/month", "$170-365/month", "$280-500+/month"],
], font_size=14)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
             "Why 60-80% Cheaper at Scale", font_size=22, color=WHITE, bold=True)

reasons = [
    "6-layer deterministic classification \u2014 avoids LLM calls",
    "80% of data processes on smallest model (Gemma3:4b)",
    "5-tier smart routing \u2014 only complex content escalates",
    "Local embeddings (embeddinggemma) \u2014 no API costs",
    "No per-seat licensing \u2014 one deployment serves all",
    "Fixed infrastructure cost, not per-query",
]
add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(2.2),
                reasons, font_size=13, color=LIGHT_GRAY)

# Hardware recs
make_table(slide, Inches(6.8), Inches(4.5), Inches(6.0), [
    ["Mac Mini", "RAM", "Monthly Cost", "What You Get"],
    ["M2 16GB", "16 GB", "$0", "Full stack, cloud LLMs"],
    ["M4 Pro 48GB", "48 GB", "$0", "Everything + local 4b + embeddings"],
    ["M4 Pro/Max 64GB", "64 GB", "$0", "All Ollama models incl. 27b"],
], font_size=12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — Competitive Comparison
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "How Mem-Dog Compares", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

make_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), [
    ["Feature", "Mem-Dog", "Mem0", "Zep", "BerryDB", "Snowflake", "Databricks"],
    ["Self-hosted / air-gapped", "\u2705", "\u274c", "\u274c", "\u274c", "\u274c", "\u274c"],
    ["Local AI models ($0)", "\u2705", "\u274c", "\u274c", "\u274c", "\u274c", "\u274c"],
    ["33 channel adapters", "\u2705", "\u274c", "\u274c", "\u274c", "\u274c", "\u274c"],
    ["300+ app integrations", "\u2705", "\u274c", "\u274c", "\u274c", "Partial", "Partial"],
    ["42 AI agents", "\u2705", "\u274c", "\u274c", "\u274c", "\u274c", "\u274c"],
    ["Temporal knowledge graph", "\u2705", "\u274c", "\u2705", "\u2705", "\u274c", "\u274c"],
    ["5 search + 4 rerankers", "\u2705", "Partial", "Partial", "Partial", "Partial", "\u274c"],
    ["RAG chat with citations", "\u2705", "Partial", "Partial", "\u2705", "\u2705", "Custom"],
    ["10 memory types + TTL", "\u2705", "Partial", "\u274c", "\u274c", "\u274c", "\u274c"],
    ["DigiMe agent (25+ apps)", "\u2705", "\u274c", "\u274c", "\u274c", "\u274c", "\u274c"],
    ["AI Studio UI", "\u2705", "\u274c", "\u274c", "Partial", "\u274c", "Partial"],
    ["MCP server", "\u2705", "\u274c", "\u274c", "\u274c", "\u274c", "\u274c"],
    ["Multi-language SDKs (5)", "\u2705", "Partial", "\u274c", "Partial", "\u2705", "\u2705"],
    ["Petabyte analytics", "\u274c", "\u274c", "\u274c", "\u274c", "\u2705", "\u2705"],
], font_size=11)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 19 — Deployment Options
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Deployment Options", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

make_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), [
    ["Option", "What You Get", "Best For"],
    ["Local (Docker Compose)", "docker compose up \u2014 11 services, hot-reload, local storage", "Dev, eval, personal use"],
    ["Google Cloud (GKE)", "Production: API + pipeline + gateway on GKE, UI on Cloud Run", "Teams, production"],
    ["Mac Mini Home Server", "Full stack on Apple Silicon. Native GPU for Ollama.", "Privacy-first, small teams"],
], font_size=14)

# Namespace layout
add_text_box(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(0.5),
             "GKE Namespace Layout", font_size=22, color=WHITE, bold=True)

make_table(slide, Inches(0.5), Inches(3.8), Inches(6.0), [
    ["Namespace", "Services"],
    ["memdog", "API, MCP Server"],
    ["webhook-pipeline", "NATS worker, 42 agents, Ollama (3 tiers)"],
    ["webhook-gateway", "Gateway, DigiMe (OpenClaw Node)"],
    ["supabase", "Postgres 16, GoTrue, Kong, PostgREST"],
    ["neo4j", "Neo4j 5.26 Community"],
    ["nango", "Nango Server + Nango Postgres"],
], font_size=12)

# Deploy commands
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.2), Inches(6.0), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
shape.line.fill.background()

add_text_box(slide, Inches(7.0), Inches(3.3), Inches(5.5), Inches(0.5),
             "Deploy Commands", font_size=18, color=ACCENT, bold=True)

deploy_code = (
    "# Local \u2014 everything\n"
    "docker compose up\n\n"
    "# GKE \u2014 per-component\n"
    "GKE_CLUSTER=open-jaw GKE_ZONE=us-central1-a \\\n"
    "  ./scripts/manual-deploy.sh deploy-api-gke \\\n"
    "  -p memdog-dev -e dev\n\n"
    "# UI \u2192 Cloud Run\n"
    "./scripts/manual-deploy.sh deploy-ui \\\n"
    "  -p memdog-dev -e dev"
)
add_text_box(slide, Inches(7.0), Inches(3.9), Inches(5.5), Inches(2.5),
             deploy_code, font_size=11, color=GREEN, font_name="Courier New")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 20 — By the Numbers
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "By the Numbers", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

stats = [
    ("42", "AI Agents"),
    ("33", "Channel Adapters"),
    ("300+", "App Integrations"),
    ("5", "Search Modes"),
    ("4", "Rerankers"),
    ("10", "Memory Types"),
    ("16+", "AI Providers"),
    ("5", "Language SDKs"),
    ("33", "API Routers"),
    ("80+", "REST Endpoints"),
    ("8", "MCP Tools"),
    ("8", "Entity Types"),
    ("6+1", "Classification Layers"),
    ("5", "Model Tiers"),
    ("11", "Docker Services"),
    ("6", "K8s Namespaces"),
]

cols = 4
rows_count = (len(stats) + cols - 1) // cols
for idx, (number, label) in enumerate(stats):
    r = idx // cols
    c = idx % cols
    left = Inches(0.5) + c * Inches(3.2)
    top = Inches(1.5) + r * Inches(1.3)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.9), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TABLE_HEADER_BG
    shape.line.fill.background()

    add_text_box(slide, left + Inches(0.1), top + Inches(0.05), Inches(2.7), Inches(0.6),
                 number, font_size=32, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.6), Inches(2.7), Inches(0.4),
                 label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Cost callout
add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
             "$0/month self-hosted  \u00b7  60-80% cheaper than cloud  \u00b7  ~30 second startup",
             font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 21 — Summary / Closing
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
             "The Private AI System", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5), Inches(1.7), Inches(3.333))

capabilities = [
    ("Ingest", "from 33 channel adapters and 300+ apps"),
    ("Enrich", "with 42 specialized AI agents across 60+ data types"),
    ("Store", "with 10 memory types, versioning, TTL, and access control"),
    ("Graph", "entities in a dual-layer temporal knowledge graph"),
    ("Search", "with 5 modes, 4 rerankers, and temporal queries"),
    ("Query", "with RAG chat and inline citations"),
    ("Manage", "models, routing, agents, and infra from AI Studio"),
    ("Connect", "via 5 SDKs, 8 MCP tools, and 80+ REST endpoints"),
    ("Monitor", "with OpenTelemetry tracing and insights dashboards"),
    ("Converse", "through DigiMe in 25+ messaging platforms"),
]

for j, (verb, rest) in enumerate(capabilities):
    top = Inches(2.2) + j * Inches(0.4)
    tf = add_text_box(slide, Inches(2.5), top, Inches(8.3), Inches(0.4),
                      "", font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT)
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{verb}  "
    run1.font.size = Pt(16)
    run1.font.color.rgb = ACCENT
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = rest
    run2.font.size = Pt(16)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = "Calibri"

add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.6),
             "Your AI. Your hardware. Your rules.",
             font_size=28, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
output_path = "/Users/paragagarwal/repos/memdog/docs/presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
