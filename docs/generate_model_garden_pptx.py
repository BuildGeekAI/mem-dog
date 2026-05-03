#!/usr/bin/env python3
"""Generate Model Garden presentation. One idea per slide, clean layout."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_HEADER_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_ROW_BG = RGBColor(0x15, 0x1F, 0x32)
TABLE_ALT_BG = RGBColor(0x1A, 0x25, 0x38)
GREEN = RGBColor(0x34, 0xD3, 0x99)
ORANGE = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG_DARK

def txt(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, f="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = f; p.alignment = a
    return tf

def title(slide, text):
    txt(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8), text, 38, WHITE, True)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.35), Inches(2.5), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

def card(slide, l, t, w, h, bc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG
    if bc: s.line.color.rgb = bc; s.line.width = Pt(1.5)
    else: s.line.fill.background()

def tbl(slide, l, t, w, rows, fs=15):
    nr, nc = len(rows), len(rows[0])
    sh = slide.shapes.add_table(nr, nc, l, t, w, Inches(0.45*nr))
    for r, row in enumerate(rows):
        for c, v in enumerate(row):
            cell = sh.table.cell(r, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(fs); p.font.name = "Calibri"
                p.font.bold = r == 0; p.font.color.rgb = ACCENT if r == 0 else WHITE
            fl = cell.fill; fl.solid()
            fl.fore_color.rgb = TABLE_HEADER_BG if r == 0 else TABLE_ALT_BG if r % 2 == 0 else TABLE_ROW_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def ns():
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); return s

# ═══════════════════════════════════════════════════════════════════════════
# 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
txt(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
    "Model Garden", 56, WHITE, True, PP_ALIGN.CENTER)
s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.0), Inches(3.333), Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
txt(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
    "How Mem-Dog Manages AI Providers, Routes Models, and Runs Inference",
    26, ACCENT, False, PP_ALIGN.CENTER)
txt(slide, Inches(2), Inches(5.2), Inches(9.333), Inches(0.8),
    "16+ providers  \u00b7  135+ models  \u00b7  per-user config  \u00b7  encrypted API keys\n"
    "5-tier smart routing  \u00b7  3-level fallback chains  \u00b7  K8s pod management",
    16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# 2 — What is Model Garden?
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "What is Model Garden?")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.0),
    "A unified provider management system where each user configures "
    "their own AI providers with encrypted API keys. The system then "
    "automatically routes each data type to the best available model.",
    20, LIGHT_GRAY)

items = [
    ("Provider Layer", ACCENT,
     "Connect 16+ AI providers.\nEncrypted API key storage.\nTest connectivity and\ndiscover available models."),
    ("Routing Layer", ORANGE,
     "Map each of 42 data types\nto the best model.\nPrimary + fallback chains.\nUser overrides supported."),
    ("Inference Layer", GREEN,
     "Unified calls via LiteLLM.\n3-level fallback on failure.\nToken usage tracking.\nPer-user credential resolution."),
]
for i, (t, c, d) in enumerate(items):
    left = Inches(1) + i * Inches(3.9)
    card(slide, left, Inches(3.2), Inches(3.5), Inches(3.2), bc=c)
    txt(slide, left+Inches(0.4), Inches(3.5), Inches(2.7), Inches(0.5), t, 22, c, True)
    txt(slide, left+Inches(0.4), Inches(4.2), Inches(2.7), Inches(1.8), d, 16, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# 3 — Supported Providers
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "16+ AI Providers")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
    "Users connect any combination of providers. Each connection is tested, "
    "models discovered, and API keys encrypted with Fernet AES-256.",
    18, LIGHT_GRAY)

tbl(slide, Inches(1), Inches(2.6), Inches(5.3), [
    ["Provider", "Type"],
    ["Ollama (Local)", "Self-hosted, free"],
    ["Ollama Cloud", "Managed Ollama"],
    ["Google Gemini", "Cloud API"],
    ["OpenAI", "Cloud API"],
    ["Anthropic", "Cloud API"],
    ["Groq", "Ultra-fast inference"],
    ["Mistral", "Cloud API"],
    ["Cohere", "Cloud API"],
], 15)

tbl(slide, Inches(7), Inches(2.6), Inches(5.333), [
    ["Provider", "Type"],
    ["DeepSeek", "Cloud API"],
    ["xAI (Grok)", "Cloud API"],
    ["OpenRouter", "Aggregator (100+ models)"],
    ["Together AI", "Cloud hosting"],
    ["HuggingFace", "Cloud / local"],
    ["AWS Bedrock", "Enterprise cloud"],
    ["vLLM", "Self-hosted serving"],
    ["LiteLLM Gateway", "Unified proxy"],
], 15)

# ═══════════════════════════════════════════════════════════════════════════
# 4 — Per-User Engine Config
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "Per-User Provider Configuration")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "Each user owns their provider connections. API keys are encrypted "
    "at rest and never returned in API responses \u2014 only a has_api_key boolean.",
    18, LIGHT_GRAY)

card(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(3.8), bc=ACCENT)
txt(slide, Inches(1.4), Inches(3.0), Inches(4.7), Inches(0.5),
    "What\u2019s Stored Per Engine", 22, ACCENT, True)
fields = [
    "Engine type (openai, gemini, ollama...)",
    "Display name and base URL",
    "API key (Fernet AES-256 encrypted)",
    "Discovered models list",
    "Available capabilities (embed, complete)",
    "Last test status + latency",
    "Enabled / disabled toggle",
]
for j, f in enumerate(fields):
    txt(slide, Inches(1.4), Inches(3.7)+j*Inches(0.38), Inches(4.7), Inches(0.38),
        f"\u2022  {f}", 16, LIGHT_GRAY)

card(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.8), bc=GREEN)
txt(slide, Inches(7.4), Inches(3.0), Inches(4.7), Inches(0.5),
    "What Users Can Do", 22, GREEN, True)
actions = [
    ("Add provider", "Connect with API key + base URL"),
    ("Test connection", "Verify reachability + latency"),
    ("Discover models", "Auto-fetch available model list"),
    ("Enable / disable", "Toggle providers on and off"),
    ("Override routing", "Pick specific models per data type"),
]
for j, (act, desc) in enumerate(actions):
    y = Inches(3.7)+j*Inches(0.5)
    txt(slide, Inches(7.4), y, Inches(2.2), Inches(0.3), act, 16, ACCENT, True)
    txt(slide, Inches(7.4), y+Inches(0.25), Inches(4.7), Inches(0.3), desc, 14, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# 5 — 5-Tier Smart Routing
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "5-Tier Smart Routing")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "Each of the 42 agent types declares a minimum model tier. "
    "The routing system picks the smallest model that can handle the job.",
    18, LIGHT_GRAY)

tbl(slide, Inches(1), Inches(2.7), Inches(11.333), [
    ["Tier", "Default Model", "Used For", "Cost"],
    ["Small", "Gemma3:4b", "JSON, CSV, YAML, XML, IoT, GPS, sensor data", "Free (local)"],
    ["Medium", "Gemma3:12b", "Code, email, chat, financial, markdown", "Free (local)"],
    ["Large", "Gemma3:27b", "PDFs, Office docs, web pages, complex documents", "Free (local)"],
    ["Multimodal", "Qwen3-VL", "Images, visual PDFs, screenshots, medical imaging", "Free (local)"],
    ["Omni", "Qwen3.5", "Audio, video, multi-format content", "Free (local)"],
], 16)

txt(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(0.8),
    "80% of data processes on the small tier (4b model). "
    "Only complex content like PDFs and images escalate to larger models.",
    18, GREEN, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# 6 — Data Type to Model Mapping
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "Data Type \u2192 Model Mapping")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "39 data types are mapped to tiers based on complexity. "
    "Users can override any mapping from the Routing tab in AI Studio.",
    18, LIGHT_GRAY)

tbl(slide, Inches(1), Inches(2.6), Inches(5.3), [
    ["Data Type", "Tier"],
    ["json, csv, yaml, xml", "Small"],
    ["sensor, gps, iot_sensor", "Small"],
    ["biometric, time_series", "Small"],
    ["code, log_file, log_stream", "Medium"],
    ["email, chat, feed", "Medium"],
    ["markdown, calendar", "Medium"],
    ["financial, channel_message", "Medium"],
], 15)

tbl(slide, Inches(7), Inches(2.6), Inches(5.333), [
    ["Data Type", "Tier"],
    ["pdf, office_doc, html_doc", "Large"],
    ["web_page, archive", "Large"],
    ["geospatial, lidar, model_3d", "Large"],
    ["scientific, infrastructure", "Large"],
    ["image, image_batch", "Multimodal"],
    ["medical_imaging", "Multimodal"],
    ["video_url, audio_stream", "Omni"],
], 15)

card(slide, Inches(1), Inches(6.0), Inches(11.333), Inches(0.9))
txt(slide, Inches(1.4), Inches(6.15), Inches(10.5), Inches(0.5),
    "Override example:  User sets pdf \u2192 openai/gpt-4o instead of Gemma3:27b. "
    "The system uses their OpenAI key for PDFs, keeps defaults for everything else.",
    15, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# 7 — Fallback Chains
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "3-Level Fallback Chains")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "If a model is unavailable, the system gracefully falls back. "
    "Fallbacks apply at the model, credential, and infrastructure levels.",
    18, LIGHT_GRAY)

chains = [
    ("Model Fallback", ACCENT,
     "When the primary model fails:",
     ["Try primary model (e.g. ollama/gemma3:12b)",
      "Try fallback model (e.g. gemini/gemini-2.5-flash)",
      "Try self-hosted model server (in-cluster)",
      "Raise error if all unavailable"]),
    ("Credential Fallback", GREEN,
     "When resolving API keys:",
     ["Check user-configured engine credentials",
      "Check environment variable (OLLAMA_CLOUD_API_KEY)",
      "Proceed without auth (local/public endpoints)"]),
    ("Infrastructure Fallback", ORANGE,
     "When choosing model server URLs:",
     ["Try tier-specific URL (MODEL_SERVER_URL_MEDIUM)",
      "Try adjacent tier URLs (small \u2192 large)",
      "Fall back to LiteLLM cloud call"]),
]
for i, (t, c, sub, steps) in enumerate(chains):
    left = Inches(0.7) + i * Inches(4.2)
    card(slide, left, Inches(2.7), Inches(3.9), Inches(4.0), bc=c)
    txt(slide, left+Inches(0.3), Inches(2.9), Inches(3.3), Inches(0.5), t, 22, c, True)
    txt(slide, left+Inches(0.3), Inches(3.4), Inches(3.3), Inches(0.4), sub, 14, WHITE)
    for j, step in enumerate(steps):
        txt(slide, left+Inches(0.3), Inches(4.0)+j*Inches(0.45), Inches(3.3), Inches(0.45),
            f"{j+1}.  {step}", 14, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# 8 — LiteLLM Unification
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "LiteLLM: One Interface, All Providers")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "All 16+ providers are called through a single LiteLLM interface. "
    "The model prefix determines which provider handles the request.",
    18, LIGHT_GRAY)

tbl(slide, Inches(1), Inches(2.7), Inches(5.5), [
    ["Prefix", "Routes To"],
    ["ollama/", "Local or Cloud Ollama"],
    ["gemini/", "Google Gemini API"],
    ["openai/", "OpenAI API"],
    ["anthropic/", "Anthropic API"],
    ["openrouter/", "OpenRouter gateway"],
    ["together_ai/", "Together AI"],
    ["huggingface/", "HuggingFace Inference"],
    ["bedrock/", "AWS Bedrock"],
], 15)

card(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(3.5), bc=ACCENT)
txt(slide, Inches(7.4), Inches(2.9), Inches(4.7), Inches(0.5),
    "How It Works", 22, ACCENT, True)

steps = [
    "Agent calls chat_for_data_type(agent_type=\"pdf\")",
    "Fetch routing config \u2192 primary: ollama/gemma3:27b",
    "Resolve credentials (user engine or env var)",
    "Call litellm.completion(model=\"ollama/gemma3:27b\")",
    "LiteLLM routes to Ollama Cloud API",
    "Track token usage (prompt + completion)",
    "Return generated text to agent",
]
for j, step in enumerate(steps):
    txt(slide, Inches(7.4), Inches(3.6)+j*Inches(0.35), Inches(4.7), Inches(0.35),
        f"{j+1}.  {step}", 14, GREEN if j in (3,4) else WHITE)

# ═══════════════════════════════════════════════════════════════════════════
# 9 — API Key Security
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "API Key Security")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "Provider API keys are encrypted at rest using Fernet (AES-256). "
    "Keys are never exposed in API responses or logs.",
    18, LIGHT_GRAY)

card(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(3.5), bc=ACCENT)
txt(slide, Inches(1.4), Inches(3.0), Inches(4.7), Inches(0.5),
    "Encryption Flow", 22, ACCENT, True)
enc_steps = [
    ("User submits API key", WHITE),
    ("encrypt_api_key() \u2192 Fernet AES-256", GREEN),
    ("Store encrypted blob in engine config", WHITE),
    ("API returns has_api_key: true (never the key)", ORANGE),
    ("Decryption only for: test, discover, inference", WHITE),
]
for j, (step, c) in enumerate(enc_steps):
    txt(slide, Inches(1.4), Inches(3.7)+j*Inches(0.45), Inches(4.7), Inches(0.45),
        f"{j+1}.  {step}", 16, c)

card(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.5), bc=GREEN)
txt(slide, Inches(7.4), Inches(3.0), Inches(4.7), Inches(0.5),
    "Security Guarantees", 22, GREEN, True)
guarantees = [
    "Master key from MASTER_ENCRYPTION_KEY env var",
    "Fernet provides authenticated encryption",
    "API responses only show has_api_key boolean",
    "Decryption happens server-side, never client",
    "If encryption unavailable, warning logged",
]
for j, g in enumerate(guarantees):
    txt(slide, Inches(7.4), Inches(3.7)+j*Inches(0.45), Inches(4.7), Inches(0.45),
        f"\u2022  {g}", 15, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# 10 — K8s Pod Management
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "K8s Pod Management")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "Create, scale, and manage Ollama model deployments directly from "
    "the AI Studio Infrastructure tab. No kubectl needed.",
    18, LIGHT_GRAY)

card(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(3.5), bc=ACCENT)
txt(slide, Inches(1.4), Inches(3.0), Inches(4.7), Inches(0.5),
    "What You Can Do", 22, ACCENT, True)
actions = [
    ("Create pods", "Name, tier (small/medium/large), models to preload"),
    ("Scale", "0 to 3 replicas with one click"),
    ("Monitor", "Status badges, CPU/memory metrics, logs"),
    ("Delete", "Remove managed deployments"),
]
for j, (act, desc) in enumerate(actions):
    y = Inches(3.7)+j*Inches(0.6)
    txt(slide, Inches(1.4), y, Inches(4.7), Inches(0.3), act, 17, ACCENT, True)
    txt(slide, Inches(1.4), y+Inches(0.28), Inches(4.7), Inches(0.3), desc, 14, LIGHT_GRAY)

card(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.5), bc=GREEN)
txt(slide, Inches(7.4), Inches(3.0), Inches(4.7), Inches(0.5),
    "How It Works", 22, GREEN, True)
details = [
    "Managed pods prefixed mdl- and labeled\nmem-dog/managed-model=true",
    "K8s RBAC restricts to webhook-pipeline\nnamespace only",
    "Pre-existing pods show as \"infrastructure\"\n(read-only, visible for monitoring)",
    "Metrics via K8s metrics API\n(CPU %, memory usage)",
]
for j, d in enumerate(details):
    txt(slide, Inches(7.4), Inches(3.7)+j*Inches(0.6), Inches(4.7), Inches(0.55),
        f"\u2022  {d}", 14, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# 11 — Agent Config Overrides
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "Per-Agent Configuration")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "Beyond model routing, each of the 42 agent types can be customized "
    "with custom prompts, output schemas, and processing flags.",
    18, LIGHT_GRAY)

tbl(slide, Inches(1), Inches(2.7), Inches(11.333), [
    ["Config Field", "What It Controls", "Example"],
    ["system_prompt", "Custom instructions for the LLM", "\"Extract only financial entities\""],
    ["output_schema", "Override the output format", "Custom JSON schema for viewpoints"],
    ["model_tier", "Force a specific tier", "Set pdf agent to \"medium\" instead of \"large\""],
    ["skills", "Attach post-processing skills", "Run sentiment analysis after classification"],
    ["parameters", "Extra inference params", "temperature: 0.1, max_tokens: 512"],
    ["enabled / disabled", "Turn AI processing on/off per type", "Disable enrichment for CSV files"],
], 15)

txt(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(0.6),
    "Configs can be set per-user or as system defaults. "
    "User configs take priority via the resolve endpoint.",
    16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# 12 — End-to-End Flow
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
title(slide, "End-to-End: PDF Arrives")
txt(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
    "How Model Garden, Smart Routing, and Fallback Chains work together.",
    18, LIGHT_GRAY)

card(slide, Inches(1), Inches(2.6), Inches(11.333), Inches(4.2))

steps = [
    ("A PDF is ingested via POST /api/v1/data", WHITE, ""),
    ("Pipeline classifies it: MIME registry \u2192 pdf agent", WHITE, "Classification"),
    ("Router looks up routing config for agent_type=\"pdf\"", WHITE, "Smart Routing"),
    ("Config says: primary=ollama/gemma3:27b, fallback=gemini/gemini-2.5-flash", ACCENT, ""),
    ("Resolve credentials: user has an Ollama Cloud engine \u2192 use their API key", WHITE, "Credentials"),
    ("Call litellm.completion(model=\"ollama/gemma3:27b\", api_key=decrypted_key)", GREEN, "Inference"),
    ("Ollama Cloud returns analysis \u2192 viewpoint + embeddings + entities stored", GREEN, ""),
    ("If Ollama fails \u2192 retry with gemini/gemini-2.5-flash using system Gemini key", ORANGE, "Fallback"),
    ("If Gemini fails \u2192 try self-hosted MODEL_SERVER_URL_LARGE", ORANGE, ""),
    ("Result: PDF fully enriched and queryable via any of 5 search modes", GREEN, ""),
]
for j, (step, color, label) in enumerate(steps):
    y = Inches(2.9)+j*Inches(0.37)
    if label:
        txt(slide, Inches(1.4), y, Inches(2.0), Inches(0.35), label, 12, ACCENT2, True)
    txt(slide, Inches(3.4), y, Inches(8.5), Inches(0.35),
        f"{j+1}.  {step}", 15, color)

# ═══════════════════════════════════════════════════════════════════════════
# 13 — Summary
# ═══════════════════════════════════════════════════════════════════════════
slide = ns()
txt(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.8),
    "Model Garden", 48, WHITE, True, PP_ALIGN.CENTER)
s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.9), Inches(3.333), Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

caps = [
    ("Connect", "16+ AI providers with encrypted API key storage"),
    ("Discover", "available models automatically from any provider"),
    ("Route", "42 data types to the smallest capable model"),
    ("Override", "any routing decision per-user from the UI"),
    ("Fallback", "gracefully across models, credentials, and infrastructure"),
    ("Encrypt", "all API keys with Fernet AES-256 at rest"),
    ("Manage", "Ollama pods from the browser \u2014 create, scale, monitor"),
    ("Customize", "agent prompts, schemas, and processing flags"),
]
for j, (v, r) in enumerate(caps):
    top = Inches(2.4)+j*Inches(0.5)
    tf = txt(slide, Inches(2.5), top, Inches(8.333), Inches(0.45), "", 20)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"{v}  "; r1.font.size = Pt(20)
    r1.font.color.rgb = ACCENT; r1.font.bold = True; r1.font.name = "Calibri"
    r2 = p.add_run(); r2.text = r; r2.font.size = Pt(20)
    r2.font.color.rgb = LIGHT_GRAY; r2.font.name = "Calibri"

txt(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.6),
    "Any provider. Any model. Automatic routing. Zero lock-in.",
    26, GREEN, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
out = "/Users/paragagarwal/repos/mem-dog/docs/model-garden.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
