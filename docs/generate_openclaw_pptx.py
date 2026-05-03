#!/usr/bin/env python3
"""Generate OpenClaw + mem-dog integration presentation as .pptx.

Design principle: one idea per slide, large text, plenty of whitespace.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_HEADER_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_ROW_BG = RGBColor(0x15, 0x1F, 0x32)
TABLE_ALT_BG = RGBColor(0x1A, 0x25, 0x38)
GREEN = RGBColor(0x34, 0xD3, 0x99)
ORANGE = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def slide_title(slide, title):
    add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
             title, size=38, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(1.35), Inches(2.5), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def card(slide, left, top, width, height, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()


def make_table(slide, left, top, width, rows_data, font_size=14):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                   width, Inches(0.45 * n_rows))
    table = shape.table
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Calibri"
                p.font.bold = r == 0
                p.font.color.rgb = ACCENT if r == 0 else WHITE
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = (TABLE_HEADER_BG if r == 0
                                   else TABLE_ALT_BG if r % 2 == 0
                                   else TABLE_ROW_BG)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
         "OpenClaw + Mem-Dog", size=60, bold=True, align=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(5), Inches(3.2), Inches(3.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
         "Turning 25+ Messaging Apps into Your AI Knowledge Base",
         size=26, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(5.5), Inches(9.333), Inches(0.5),
         "WhatsApp  \u00b7  Signal  \u00b7  Telegram  \u00b7  Slack  \u00b7  Discord  "
         "\u00b7  Matrix  \u00b7  MS Teams  \u00b7  25+ more",
         size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is DigiMe?
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "What is DigiMe?")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
         "An AI assistant built on the OpenClaw runtime that lives inside "
         "your messaging apps. Every message is captured, enriched by 42 AI "
         "agents, and stored in a temporal knowledge graph \u2014 so you can "
         "ask questions across all your conversations from any channel.",
         size=20, color=LIGHT_GRAY)

items = [
    ("OpenClaw Runtime", ACCENT,
     "Open-source multi-channel agent framework.\n"
     "Handles channel connections, WebSocket\n"
     "persistence, and Gemini LLM orchestration."),
    ("Mem-Dog Backend", GREEN,
     "42 AI agents classify and enrich every\n"
     "message. Dual-layer knowledge graph.\n"
     "5 search modes with RAG chat."),
    ("4 Bridge Skills", ORANGE,
     "Bridge (record all), Ingest (store),\n"
     "Query (RAG chat), Semantic Search\n"
     "(vector similarity)."),
]

for i, (title, color, desc) in enumerate(items):
    left = Inches(1) + i * Inches(3.9)
    card(slide, left, Inches(3.4), Inches(3.5), Inches(3.0), border_color=color)
    add_text(slide, left + Inches(0.4), Inches(3.7), Inches(2.7), Inches(0.5),
             title, size=22, color=color, bold=True)
    add_text(slide, left + Inches(0.4), Inches(4.4), Inches(2.7), Inches(1.6),
             desc, size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture (high-level flow only)
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "How It Works")

# Single horizontal flow — large boxes with arrows
flow = [
    ("25+ Channels", "WhatsApp, Signal,\nSlack, Telegram...", ACCENT),
    ("OpenClaw", "Channel adapters\nGemini LLM agent", ACCENT2),
    ("4 Skills", "Bridge, Ingest,\nQuery, Search", ORANGE),
    ("Webhook\nGateway", "Normalize + resolve\nuser identity", GREEN),
    ("Mem-Dog API\n+ Pipeline", "Store, classify,\nenrich (42 agents)", ACCENT),
    ("Knowledge\nGraph + RAG", "Search, cite,\nrespond", GREEN),
]

box_w = Inches(1.7)
gap = Inches(0.25)
total = len(flow) * box_w + (len(flow) - 1) * gap
start_x = (Inches(13.333) - total) / 2

for i, (title, sub, color) in enumerate(flow):
    left = start_x + i * (box_w + gap)
    card(slide, left, Inches(2.3), box_w, Inches(2.2), border_color=color)
    add_text(slide, left + Inches(0.15), Inches(2.5), box_w - Inches(0.3), Inches(0.8),
             title, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(3.3), box_w - Inches(0.3), Inches(0.9),
             sub, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        arrow_x = left + box_w + Inches(0.02)
        add_text(slide, arrow_x, Inches(3.0), gap, Inches(0.5),
                 "\u2192", size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.6),
         "Every message flows end-to-end: captured, classified, enriched, "
         "embedded, graphed, and made searchable \u2014 automatically.",
         size=18, color=GREEN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Bridge Skill
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "Skill 1: mem-dog-bridge")

add_text(slide, Inches(1), Inches(1.8), Inches(5), Inches(0.5),
         "PASSIVE  \u2014  fires on every message, always", size=16, color=ACCENT)

add_text(slide, Inches(1), Inches(2.5), Inches(5), Inches(2.5),
         "Every incoming message is automatically forwarded to the "
         "mem-dog webhook pipeline. Non-negotiable.\n\n"
         "The pipeline runs 6-layer classification, routes to one of "
         "42 specialized agents, generates a viewpoint + embedding, "
         "extracts entities, and dual-writes to the knowledge graph.",
         size=18, color=LIGHT_GRAY)

card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(3.0))
add_text(slide, Inches(7.4), Inches(2.0), Inches(4.7), Inches(0.5),
         "Flow", size=20, color=ACCENT, bold=True)

steps = [
    "User sends message on any channel",
    "OpenClaw \u2192 DigiMe agent",
    "Bridge skill fires automatically",
    "POST /webhooks/openclaw \u2192 Gateway",
    "Gateway normalizes \u2192 UniversalEnvelope",
    "API stores + publishes to NATS",
    "Pipeline: classify \u2192 agent \u2192 enrich \u2192 graph",
]
for j, step in enumerate(steps):
    add_text(slide, Inches(7.4), Inches(2.6) + j * Inches(0.35),
             Inches(4.7), Inches(0.35),
             f"{j+1}.  {step}", size=15,
             color=GREEN if j in (2, 6) else WHITE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Active Skills (ingest, query, search)
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "Skills 2-4: Ingest, Query, Search")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
         "ACTIVE  \u2014  triggered by user intent (questions, commands, search requests)",
         size=16, color=ORANGE)

skills = [
    ("mem-dog-ingest", "Store Data", ORANGE,
     "User says \"remember my flight is at 3pm\"\nor uses /save command.",
     "POST /api/v1/data"),
    ("mem-dog-query", "RAG Chat", GREEN,
     "User asks a question. DigiMe retrieves\ncontext and answers with [1][2] citations.",
     "POST /api/v1/ai/query"),
    ("mem-dog-semantic-search", "Vector Search", ACCENT2,
     "Raw similarity search without LLM synthesis.\nReturns ranked results with scores.",
     "POST /api/v1/ai/query/semantic"),
]

for i, (name, title, color, desc, endpoint) in enumerate(skills):
    top = Inches(2.6) + i * Inches(1.6)
    card(slide, Inches(1), top, Inches(11.333), Inches(1.3), border_color=color)

    add_text(slide, Inches(1.4), top + Inches(0.15), Inches(2.8), Inches(0.4),
             name, size=18, color=color, bold=True)
    add_text(slide, Inches(1.4), top + Inches(0.55), Inches(2.8), Inches(0.6),
             title, size=22, color=WHITE, bold=True)

    add_text(slide, Inches(4.8), top + Inches(0.25), Inches(4.5), Inches(0.8),
             desc, size=16, color=LIGHT_GRAY)

    add_text(slide, Inches(9.8), top + Inches(0.35), Inches(2.3), Inches(0.5),
             endpoint, size=14, color=GREEN, font="Courier New")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Identity Resolution
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "Multi-User Identity Resolution")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.8),
         "One DigiMe instance serves multiple users with complete data isolation. "
         "Each sender is mapped to a user_id via their channel identity.",
         size=19, color=LIGHT_GRAY)

# Left: how it works
card(slide, Inches(1), Inches(3.0), Inches(5.5), Inches(3.5))
add_text(slide, Inches(1.4), Inches(3.2), Inches(4.7), Inches(0.5),
         "Resolution Flow", size=22, color=ACCENT, bold=True)

steps = [
    ("Message arrives at webhook-gateway", WHITE),
    ("Extract (channel_type, peer_id)", WHITE),
    ("resolve_user_id() looks up identity", WHITE),
    ("  \u2192  Supabase direct read (fast path)", LIGHT_GRAY),
    ("  \u2192  Fallback: API lookup + 5-min cache", LIGHT_GRAY),
    ("All API calls scoped to resolved user_id", GREEN),
]
for j, (text, color) in enumerate(steps):
    add_text(slide, Inches(1.4), Inches(3.9) + j * Inches(0.4),
             Inches(4.7), Inches(0.4), text, size=16, color=color)

# Right: example
card(slide, Inches(7), Inches(3.0), Inches(5.5), Inches(3.5))
add_text(slide, Inches(7.4), Inches(3.2), Inches(4.7), Inches(0.5),
         "Example", size=22, color=ORANGE, bold=True)

add_text(slide, Inches(7.4), Inches(3.9), Inches(4.7), Inches(2.2),
         "Alice registers 3 channel identities:\n\n"
         "  (whatsapp, +1925...)   \u2192  user_alice\n"
         "  (signal, alice_sig)    \u2192  user_alice\n"
         "  (slack, U04ABC)        \u2192  user_alice\n\n"
         "All 3 channels share one unified\n"
         "knowledge base. Bob can\u2019t see Alice\u2019s data.",
         size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Supported Channels
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "25+ Supported Channels")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "All channels routed through a single OpenClawBridgeAdapter. "
         "No per-channel code needed.",
         size=18, color=LIGHT_GRAY)

make_table(slide, Inches(1), Inches(2.7), Inches(11.333), [
    ["Category", "Channels"],
    ["Encrypted Messaging", "Signal, Matrix, WhatsApp (bridge)"],
    ["Chat Platforms", "Slack, Discord, MS Teams, IRC, Google Chat"],
    ["Asian Messaging", "Line, Feishu, Zalo"],
    ["Self-Hosted", "Mattermost, Nextcloud Talk, Rocket.Chat, Synology Chat"],
    ["Decentralized", "Nostr, Matrix, Tlon"],
    ["Apple / Streaming / Web", "BlueBubbles / iMessage, Twitch, Webchat"],
    ["Social", "WeChat, Viber"],
], font_size=16)

add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
         "Flexible field normalization: channel/channelType, sender/from, "
         "text/body/message/content \u2014 all mapped to UniversalEnvelope.",
         size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Deployment
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()
slide_title(slide, "Deployment")

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         "Single K8s pod in the webhook-gateway namespace. "
         "Serves all 25+ channels simultaneously.",
         size=18, color=LIGHT_GRAY)

make_table(slide, Inches(1), Inches(2.7), Inches(5.5), [
    ["Resource", "Details"],
    ["Namespace", "webhook-gateway (shared)"],
    ["Deployment", "1 replica, 2 CPU / 4Gi limit"],
    ["Service", "ClusterIP :18789"],
    ["HTTPRoute", "/oc/* \u2192 strips prefix"],
    ["PVC", "openclaw-home (persistent state)"],
    ["Secrets", "Gemini key, API key, gateway token"],
], font_size=15)

card(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(3.5))
add_text(slide, Inches(7.4), Inches(2.9), Inches(4.7), Inches(0.5),
         "Key Config", size=22, color=ACCENT, bold=True)

config_lines = [
    ("GEMINI_MODEL", "gemini-3.1-pro-preview"),
    ("MEM_DOG_API_URL", "http://api.mem-dog.svc:8080"),
    ("BRIDGE_URL", "http://webhook-gateway.svc:8080\n/webhooks/openclaw"),
    ("SOUL.md", "Always record via bridge,\nnever use local files"),
    ("IDENTITY.md", "DigiMe persona definition"),
]
for j, (key, val) in enumerate(config_lines):
    y = Inches(3.5) + j * Inches(0.5)
    add_text(slide, Inches(7.4), y, Inches(2.0), Inches(0.45),
             key, size=14, color=ACCENT, bold=True)
    add_text(slide, Inches(9.4), y, Inches(2.8), Inches(0.45),
             val, size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Summary
# ═══════════════════════════════════════════════════════════════════════════
slide = new_slide()

add_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.8),
         "OpenClaw + Mem-Dog", size=48, bold=True, align=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(5), Inches(1.9), Inches(3.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

capabilities = [
    ("Connect", "to 25+ messaging platforms via OpenClaw"),
    ("Record", "every message automatically via bridge skill"),
    ("Enrich", "with 42 AI agents \u2014 classify, embed, extract"),
    ("Store", "in dual-layer temporal knowledge graph"),
    ("Search", "with 5 modes + 4 rerankers"),
    ("Answer", "with RAG chat and [1][2] citations"),
    ("Isolate", "per-user data via channel identity"),
    ("Deploy", "as a single K8s pod for all channels"),
]

for j, (verb, rest) in enumerate(capabilities):
    top = Inches(2.5) + j * Inches(0.48)
    tf = add_text(slide, Inches(3), top, Inches(7.333), Inches(0.45),
                  "", size=20, align=PP_ALIGN.LEFT)
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{verb}  "
    run1.font.size = Pt(20)
    run1.font.color.rgb = ACCENT
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = rest
    run2.font.size = Pt(20)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = "Calibri"

add_text(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.6),
         "Any channel. One agent. All your knowledge.",
         size=30, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
output_path = "/Users/paragagarwal/repos/mem-dog/docs/openclaw-integration.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
