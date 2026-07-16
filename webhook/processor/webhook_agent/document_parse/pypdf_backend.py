"""pypdf / python-docx extractors — baseline parser backend."""

from __future__ import annotations

import io
import logging
import time
from typing import Any

from .config import DOCUMENT_PARSE_MAX_PAGES, PYPDF_LEGACY_MAX_PAGES
from .models import ParsedDocument

logger = logging.getLogger("mem_dog.webhook.document_parse.pypdf")

try:
    import pypdf
    _PYPDF_VERSION = getattr(pypdf, "__version__", "unknown")
except ImportError:
    _PYPDF_VERSION = "missing"


def _pdf_max_pages() -> int:
    if DOCUMENT_PARSE_MAX_PAGES > 0:
        return DOCUMENT_PARSE_MAX_PAGES
    return PYPDF_LEGACY_MAX_PAGES


def parse_pdf(content_bytes: bytes, mime_type: str) -> ParsedDocument:
    from pypdf import PdfReader

    t0 = time.perf_counter()
    reader = PdfReader(io.BytesIO(content_bytes))
    total_pages = len(reader.pages)
    max_pages = _pdf_max_pages()
    pages = reader.pages[:max_pages] if max_pages > 0 else reader.pages

    elements: list[dict[str, Any]] = []
    text_parts: list[str] = []
    pages_with_text = 0
    for i, page in enumerate(pages):
        page_text = page.extract_text() or ""
        if page_text.strip():
            pages_with_text += 1
            text_parts.append(f"[Page {i + 1}]\n{page_text}")
            elements.append(
                {
                    "type": "page",
                    "text": page_text.strip(),
                    "page": i + 1,
                    "section_path": [f"Page {i + 1}"],
                    "element_type": "narrative",
                }
            )

    markdown = "\n\n".join(text_parts)
    latency = int((time.perf_counter() - t0) * 1000)
    confidence = pages_with_text / len(pages) if pages else 0.0

    return ParsedDocument(
        parser="pypdf",
        parser_version=_PYPDF_VERSION,
        mime_type=mime_type,
        page_count=total_pages,
        ocr_used=False,
        confidence=round(confidence, 3),
        markdown=markdown,
        elements=elements,
        chunks=[
            {
                "text": el["text"],
                "page": el["page"],
                "section_path": el["section_path"],
                "element_type": el["element_type"],
            }
            for el in elements
        ],
        parse_latency_ms=latency,
    )


def parse_docx(content_bytes: bytes, mime_type: str) -> ParsedDocument:
    from docx import Document

    t0 = time.perf_counter()
    doc = Document(io.BytesIO(content_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    markdown = "\n".join(paragraphs)
    elements = [
        {
            "type": "paragraph",
            "text": t,
            "page": None,
            "section_path": [],
            "element_type": "narrative",
        }
        for t in paragraphs
    ]
    latency = int((time.perf_counter() - t0) * 1000)
    return ParsedDocument(
        parser="pypdf",
        parser_version=_PYPDF_VERSION,
        mime_type=mime_type,
        page_count=0,
        ocr_used=False,
        confidence=1.0 if markdown else 0.0,
        markdown=markdown,
        elements=elements,
        chunks=[{"text": t, "page": None, "section_path": [], "element_type": "narrative"} for t in paragraphs],
        parse_latency_ms=latency,
    )


def parse_xlsx(content_bytes: bytes, mime_type: str) -> ParsedDocument:
    from openpyxl import load_workbook

    t0 = time.perf_counter()
    wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
    parts: list[str] = []
    elements: list[dict[str, Any]] = []
    for sheet_name in wb.sheetnames[:10]:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(max_row=100, values_only=True):
            row_text = ", ".join(str(c) for c in row if c is not None)
            if row_text:
                rows.append(row_text)
        if rows:
            block = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
            parts.append(block)
            elements.append(
                {
                    "type": "sheet",
                    "text": block,
                    "page": None,
                    "section_path": [sheet_name],
                    "element_type": "table",
                }
            )
    wb.close()
    markdown = "\n\n".join(parts)
    latency = int((time.perf_counter() - t0) * 1000)
    return ParsedDocument(
        parser="pypdf",
        parser_version=_PYPDF_VERSION,
        mime_type=mime_type,
        page_count=0,
        ocr_used=False,
        confidence=1.0 if markdown else 0.0,
        markdown=markdown,
        elements=elements,
        chunks=[
            {
                "text": el["text"],
                "page": None,
                "section_path": el["section_path"],
                "element_type": el["element_type"],
            }
            for el in elements
        ],
        parse_latency_ms=latency,
    )


def parse_pptx(content_bytes: bytes, mime_type: str) -> ParsedDocument:
    from pptx import Presentation

    t0 = time.perf_counter()
    prs = Presentation(io.BytesIO(content_bytes))
    parts: list[str] = []
    elements: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            block = f"[Slide {i + 1}]\n" + "\n".join(texts)
            parts.append(block)
            elements.append(
                {
                    "type": "slide",
                    "text": block,
                    "page": i + 1,
                    "section_path": [f"Slide {i + 1}"],
                    "element_type": "narrative",
                }
            )
    markdown = "\n\n".join(parts)
    latency = int((time.perf_counter() - t0) * 1000)
    return ParsedDocument(
        parser="pypdf",
        parser_version=_PYPDF_VERSION,
        mime_type=mime_type,
        page_count=len(prs.slides),
        ocr_used=False,
        confidence=1.0 if markdown else 0.0,
        markdown=markdown,
        elements=elements,
        chunks=[
            {
                "text": el["text"],
                "page": el.get("page"),
                "section_path": el["section_path"],
                "element_type": el["element_type"],
            }
            for el in elements
        ],
        parse_latency_ms=latency,
    )


def parse(content_bytes: bytes, mime_type: str) -> ParsedDocument:
    mime_lower = (mime_type or "").lower()

    if "pdf" in mime_lower:
        return parse_pdf(content_bytes, mime_type)

    if "wordprocessingml" in mime_lower or mime_lower == "application/msword":
        if "msword" in mime_lower and "officedocument" not in mime_lower:
            msg = f"[Legacy .doc format — metadata only] File size: {len(content_bytes)} bytes"
            return ParsedDocument(
                parser="pypdf",
                parser_version=_PYPDF_VERSION,
                mime_type=mime_type,
                page_count=0,
                ocr_used=False,
                confidence=0.0,
                markdown=msg,
            )
        return parse_docx(content_bytes, mime_type)

    if "spreadsheetml" in mime_lower or "ms-excel" in mime_lower:
        if "ms-excel" in mime_lower and "officedocument" not in mime_lower:
            msg = f"[Legacy .xls format — metadata only] File size: {len(content_bytes)} bytes"
            return ParsedDocument(
                parser="pypdf",
                parser_version=_PYPDF_VERSION,
                mime_type=mime_type,
                page_count=0,
                ocr_used=False,
                confidence=0.0,
                markdown=msg,
            )
        return parse_xlsx(content_bytes, mime_type)

    if "presentationml" in mime_lower or "ms-powerpoint" in mime_lower:
        if "ms-powerpoint" in mime_lower and "officedocument" not in mime_lower:
            msg = f"[Legacy .ppt format — metadata only] File size: {len(content_bytes)} bytes"
            return ParsedDocument(
                parser="pypdf",
                parser_version=_PYPDF_VERSION,
                mime_type=mime_type,
                page_count=0,
                ocr_used=False,
                confidence=0.0,
                markdown=msg,
            )
        return parse_pptx(content_bytes, mime_type)

    return ParsedDocument(
        parser="pypdf",
        parser_version=_PYPDF_VERSION,
        mime_type=mime_type,
        page_count=0,
        ocr_used=False,
        confidence=0.0,
        markdown="",
        error=f"unsupported mime for pypdf backend: {mime_type}",
    )
